import streamlit as st
import streamlit_chat as stc  
import google.generativeai as genai
import PyPDF2
import docx
import os
import platform
import random
# Changed mysql.connector to pymongo
from pymongo import MongoClient
from pymongo.errors import ConnectionFailure, OperationFailure
from bson.objectid import ObjectId # Import ObjectId for MongoDB's _id
from werkzeug.security import generate_password_hash, check_password_hash
import pandas as pd
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import nltk
from nltk.sentiment import SentimentIntensityAnalyzer
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
import re
import json
import streamlit as st  
import pandas as pd  
import matplotlib.pyplot as plt
from mpl_toolkits.mplot3d import Axes3D  
import seaborn as sns  
import googletrans
import numpy as np 
import io
import datetime
from datetime import datetime, timedelta
from fpdf import FPDF  
import csv
import tempfile
from pptx import Presentation
from streamlit_option_menu import option_menu 
from pptx.util import Inches
import plotly.express as px
from streamlit_lottie import st_lottie  # Import the Lottie function
import requests  # To fetch the Lottie animation
import googletrans
from google_trans_new import google_translator  
from deep_translator import GoogleTranslator  # Using deep-translator as it's more reliable


# Initialize translator (from google_trans_new, if used elsewhere)
translator = google_translator()


st.set_page_config(layout="wide")


def trainer_dashboard():
     # Initialize session state variables if they do not exist
    if 'generated_questions' not in st.session_state:
        st.session_state.generated_questions = []
    if 'generated_options' not in st.session_state:
        st.session_state.generated_options = []
    if 'generated_qb_id' not in st.session_state:
        st.session_state.generated_qb_id = None  
    if 'chat_history' not in st.session_state:
        st.session_state.chat_history = []  # Initialize chat history
        
    
    with st.sidebar:
      # Create a navigation bar using option_menu
      selected = option_menu(
        menu_title="Trainer dashboard",  # required
        options=[
            "Upload Curriculum", 
            "Generate Question Bank", 
            "View Questions", 
            "Review Feedback",  
            "Curriculum Overview",  # Updated icon added below
            "Download Questions",
            "Employee Performance",
            "Generate Questions",
            "Chatbot"
        ],  # required
        icons=["upload", "question-circle", "eye", "check-circle", "book", "download","bar-chart", "magic", "chat" ],  # optional
        menu_icon="cast",  # optional
        default_index=0,  # optional
        orientation="vertical",
    )
    
     # Initialize question_banks to avoid UnboundLocalError
    question_banks = []  
    # Display content based on the selected option
    if selected == "Upload Curriculum":
        st.subheader("Upload Curriculum 📁")
        technology = st.text_input("Technology", key="upload_technology")
        topics = st.text_area("Topics (one per line)", key="upload_topics")
        uploaded_file = st.file_uploader("Upload curriculum file", type=None, key="curriculum_file")

        if st.button("Upload Curriculum", key="upload_curriculum_button"):
            topic_list = [topic.strip() for topic in topics.split('\n') if topic.strip()]
            content = ""

            if uploaded_file is not None:
                try:
                    file_content = extract_text_from_file(uploaded_file)
                    content = file_content
                    topic_list.extend([topic.strip() for topic in file_content.split('\n') if topic.strip()])
                    topic_list = list(set(topic_list))  # Remove duplicates
                except ValueError as e:
                    st.error(f"Error processing file: {str(e)}")
                    return

            if upload_curriculum(technology, topic_list, content):
                st.success("Curriculum uploaded successfully!")
            else:
                st.error("Failed to upload curriculum")

    elif selected == "Generate Question Bank":
        st.subheader("Generate Question Bank 📚")
        curricula = get_all_curricula()

        if not curricula:
            st.warning("No curricula available. Please upload a curriculum first.")
        else:
            selected_curriculum = st.selectbox("Select Curriculum", options=[c['technology'] for c in curricula], key="selected_curriculum")
            if selected_curriculum:
                qb_technology = selected_curriculum
                st.write(f"Selected Technology: {qb_technology}")
                num_questions = st.number_input("Number of Questions", min_value=1, value=10, key="num_questions")
                question_type = st.selectbox("Question Type", ["multiple-choice", "subjective", "fill-in-the-blank"], key="question_type")
                difficulty = st.selectbox("Difficulty", ["Easy", "Medium", "Hard"], key="question_difficulty")

                if st.button("Generate Question Bank", key="generate_qb_button"):
                    curriculum_content = get_curriculum_text(qb_technology)
                    if curriculum_content:
                        questions, options, correct_answers = generate_questions(curriculum_content, num_questions, question_type)

                        question_bank_id = save_question_bank(
    qb_technology,
    [],
    '\n'.join(questions),
    difficulty,
    '\n'.join(correct_answers),
    question_type=question_type,
    options='\n'.join(['###'.join(opt) for opt in options])
)


                        if question_bank_id:
                            st.success(f"Question Bank generated successfully! ID: {question_bank_id}")
                            st.session_state.generated_questions = questions
                            st.session_state.generated_options = options
                            st.session_state.generated_qb_id = question_bank_id
                        else:
                            st.error("Failed to save question bank")
                    else:
                        st.error("Failed to retrieve curriculum content")

    elif selected == "View Questions":
        st.subheader("View Questions 📖")
        question_banks = get_all_question_banks()  # Fetch all question banks

        if not question_banks:
            st.info("No question banks available yet.")
        else:
            # Create a container for search and dropdown
            search_col1, search_col2 = st.columns([1, 2])
            
            with search_col1:
                search_id = st.text_input(
                    "Search by ID",
                    key="qb_search_id",
                    placeholder="Enter ID...",
                    help="Enter a question bank ID to quickly find it"
                )
            
            # Format options for the dropdown
            dropdown_options = [(str(qb['_id']), f"{qb['technology']} - {qb['difficulty']}") for qb in question_banks]
            
            # Filter dropdown options if ID is entered
            if search_id:
                try:
                    # MongoDB uses ObjectId for _id, so convert search_id to ObjectId
                    search_object_id = ObjectId(search_id)
                    dropdown_options = [(str(qb['_id']), f"{qb['technology']} - {qb['difficulty']}")
                                        for qb in question_banks if qb['_id'] == search_object_id]
                    if not dropdown_options:
                        st.warning(f"No question bank found with ID: {search_id}")
                except Exception: # Catching general exception for ObjectId conversion errors
                    st.error("Please enter a valid ID")
            
            with search_col2:
                selected_qb = st.selectbox(
                    "Select Question Bank", 
                    options=dropdown_options,
                    format_func=lambda x: f"ID: {x[0]} - {x[1]}", 
                    key="view_qb_select",
                    help="Select a question bank from the dropdown or use the ID search to filter"
                )

            # Display question bank details if selected
            if selected_qb:
                qb_id_str, _ = selected_qb
                qb_id = ObjectId(qb_id_str) # Convert back to ObjectId for database query
                qb_details = next((qb for qb in question_banks if qb['_id'] == qb_id), None)
                
                if qb_details:
                    # Display metadata in columns
                    col1, col2, col3 = st.columns(3)  # Added an extra column for updated timestamp
                    with col1:
                        st.info(f"**ID:** {qb_details['_id']}")
                    with col2:
                        st.info(f"**Technology:** {qb_details['technology']}")
                    with col3:
                        st.info(f"**Difficulty:** {qb_details['difficulty']}")
                    
                    
                    st.write("---")  # Add a separator
                    st.subheader("Questions:")

                    # Check if these are the recently generated questions
                    if 'generated_qb_id' in st.session_state and st.session_state.generated_qb_id == str(qb_details['_id']):
                        for i, (question, options) in enumerate(zip(st.session_state.generated_questions, st.session_state.generated_options), 1):
                            st.write(f"**Question {i}:** {question}")
                            for j, option in enumerate(options):
                                st.write(f"{chr(65+j)}) {option}")
                            st.write("")  # Add space between questions
                    else:
                        # Use the questions from qb_details
                        questions = qb_details.get('questions', '').split('\n')
                        for i, question in enumerate(questions, 1):
                            if question.strip():  # Only show non-empty questions
                                st.write(f"**Question {i}:** {question}")
                                
                                # Display options if the question type is multiple-choice
                                if qb_details.get('question_type', '') == 'multiple-choice':
                                    options = qb_details.get('answer', '').split('###')
                                    if options:
                                        for j, option in enumerate(options):
                                            if option.strip():  # Only show non-empty options
                                                st.write(f"{chr(65 + j)}) {option.strip()}")
                                st.write("")  # Add space between questions

                    
    elif selected == "Review Feedback":
        st.subheader("Review Feedback 🔍")
        feedback = review_feedback()
        if not feedback:
            st.info("No feedback available yet.")
        else:
            feedback_df = pd.DataFrame(feedback)
            st.dataframe(feedback_df)

            # Debugging: Check the structure of question_banks
            question_banks = get_all_question_banks()  # Ensure question_banks is populated
            st.write("Question Banks DataFrame:")
            st.write(pd.DataFrame(question_banks))  # This will show the structure of the DataFrame

            # Check if 'id' exists in question_banks (or _id for MongoDB)
            if not question_banks or '_id' not in question_banks[0]:
                st.error("The '_id' column is missing from question_banks.")
                return

            # Sentiment Analysis Summary
            sentiment_counts = feedback_df['sentiment'].value_counts()
            st.subheader("Sentiment Analysis Summary")
            st.bar_chart(sentiment_counts)

            # Question Bank Feedback Summary
            st.subheader("Question Bank Feedback Summary")
            # Convert ObjectId to string for merging
            feedback_df['question_bank_id_str'] = feedback_df['question_bank_id'].apply(str)
            qb_df = pd.DataFrame(question_banks)
            qb_df['id_str'] = qb_df['_id'].apply(str)

            qb_feedback = feedback_df.groupby('question_bank_id_str')['rating'].mean().reset_index()
            qb_feedback = qb_feedback.merge(qb_df, left_on='question_bank_id_str', right_on='id_str')
            st.dataframe(qb_feedback[['question_bank_id_str', 'technology', 'difficulty', 'rating']]) 

    elif selected == "Curriculum Overview":
        st.subheader("Curriculum Overview 📜")
        curricula = get_all_curricula()
        
        if curricula:
            # Convert ASCII values to strings for each curriculum (if applicable)
            for curriculum in curricula:
                # MongoDB stores topics as strings or lists directly, so no ASCII conversion needed unless explicitly stored as ASCII values
                if isinstance(curriculum.get('topics'), list) and all(isinstance(x, int) for x in curriculum['topics']):
                    curriculum['topics'] = ascii_to_string(curriculum['topics'])  # Convert to string

            curriculum_df = pd.DataFrame(curricula)
            st.dataframe(curriculum_df[['technology']])
        else:
            st.info("No curricula available.")

        

        # Display feedback summary for question banks
        st.subheader("Question Bank Feedback Summary")
        feedback = review_feedback()  # Assuming this function retrieves feedback data

        if feedback:
            feedback_df = pd.DataFrame(feedback)
            st.dataframe(feedback_df)  # Display feedback DataFrame
        else:
            st.info("No feedback available yet.")

        st.subheader("Clear History")  
        if st.button("Clear Curriculum Content History"):   
            db = create_connection()  
            if db is not None:  
                try:
                    db.curriculum.delete_many({}) # Clear the curriculum collection
                    st.success("Curriculum content history cleared successfully!")  
                except OperationFailure as e:
                    st.error(f"Failed to clear curriculum history: {e}")
            else:  
                st.error("Failed to connect to database")  

        if st.button("Clear Generated Topics History"):  
            db = create_connection()  
            if db is not None:  
                try:
                    db.generated_question_files.delete_many({}) # Clear the generated_question_files collection
                    st.success("Generated topics history cleared successfully!")  
                except OperationFailure as e:
                    st.error(f"Failed to clear generated topics history: {e}")
            else:  
                st.error("Failed to connect to database")     

    elif selected == "Download Questions":
        st.subheader("Download Questions ⬇️")  
        question_banks = get_all_question_banks()  
        if not question_banks:  
            st.info("No question banks available yet.")  
        else:  
            selected_qb = st.selectbox("Select Question Bank", options=[(str(qb['_id']), f"{qb['technology']} - {qb['difficulty']}") for qb in question_banks], format_func=lambda x: f"ID: {x[0]} - {x[1]}", key="download_qb_select")  
            if selected_qb:  
                qb_id_str, _ = selected_qb  
                qb_id = ObjectId(qb_id_str) # Convert to ObjectId
                qb_details = next((qb for qb in question_banks if qb['_id'] == qb_id), None)  
                if qb_details:  
                    st.write(f"Technology: {qb_details['technology']}")  
                    st.write(f"Difficulty: {qb_details['difficulty']}")  

                    file_format = st.selectbox("Select File Format", ["docx", "pdf", "pptx", "csv"])  
                    questions = qb_details['questions'].split('\n')  
                    
                    if file_format == "pdf":  
                        pdf = FPDF()  
                        pdf.add_page()  
                        pdf.set_font("Arial", size=15)  
                        for question in questions:  
                            pdf.cell(200, 10, txt=question, ln=True, align='L')  

                        # Use a temporary file to save the PDF
                        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
                            pdf.output(temp_file.name)
                            temp_file.seek(0)  # Go back to the start of the file

                            # Download button
                            st.download_button(label='Download PDF', data=temp_file.read(), file_name=f'questions_{qb_id_str}.pdf', mime='application/pdf')

                    elif file_format == "docx":  
                        doc = docx.Document()  
                        for question in questions:  
                            doc.add_paragraph(question)  
                        buffer = io.BytesIO()  
                        doc.save(buffer)  
                        buffer.seek(0)  

                        st.download_button(label='Download DOCX', data=buffer.getvalue(), file_name=f'questions_{qb_id_str}.docx', mime='application/vnd.openxmlformats-officedocument.wordprocessingml.document')  
                    
                    elif file_format == "pptx":
                        presentation = Presentation()
                        slide_layout = presentation.slide_layouts[6]  # Blank slide layout
                        slide = presentation.slides.add_slide(slide_layout)
                        
                        # Set the starting position for the textboxes in inches
                        left = Inches(1)  # 1 inch from the left
                        top = Inches(1)   # 1 inch from the top
                        height = Inches(0.5)  # Height of each textbox (0.5 inch)
                        
                        for question in questions:
                            # Add a textbox with specified dimensions
                            textbox = slide.shapes.add_textbox(left, top, width=Inches(8), height=height)  # Width of 8 inches
                            textbox.text = question
                            top += height + Inches(0.1)  # Move down for the next textbox (0.1 inch gap)
                        
                        buffer = io.BytesIO()
                        presentation.save(buffer)
                        buffer.seek(0)

                        st.download_button(label='Download PPTX', data=buffer.getvalue(), file_name=f'questions_{qb_id_str}.pptx', mime='application/vnd.openxmlformats-officedocument.presentationml.presentation')
                    
                    

                    elif file_format == "csv":
                        # Create a buffer to hold the CSV data
                        buffer = io.StringIO()
                        
                        # Create a CSV writer object
                        csv_writer = csv.writer(buffer)
                        
                        # Optionally write a header (if you want to include column names)
                        csv_writer.writerow(["Question"])  # Header for the CSV
                        
                        # Write each question to a new row in the CSV
                        for question in questions:
                            csv_writer.writerow([question])  # Each question in a new row
                        
                        # Get the CSV content from the buffer
                        buffer.seek(0)  # Go back to the start of the StringIO buffer

                        # Provide the download button for the CSV file
                        st.download_button(label='Download CSV', data=buffer.getvalue(), file_name=f'questions_{qb_id_str}.csv', mime='text/csv')

                    

    # elif selected == "Generate Questions from Prompt":
    #     st.subheader("Generate Questions from Prompt")    
    #     topic_name = st.text_input("Enter Topic Name")    
    #     prompt = st.text_area("Enter a paragraph to generate questions")    
    #     question_type = st.selectbox("Select Question Type", ["Multiple Choice", "Subjective", "Fill in the Blanks"])    
    #     difficulty_level = st.selectbox("Select Difficulty Level", ["Easy", "Medium", "Hard"])  
    #     num_questions = st.number_input("Number of Questions to Generate", min_value=1, value=10)  
        
    #     if st.button("Generate Questions"):    
    #         generated_questions = generate_questions_from_prompt(prompt, question_type, difficulty_level, num_questions, topic_name)    
    #         if generated_questions:    
    #             st.write("Generated Questions:")    
    #             selected_questions = []
    #             for i, question in enumerate(generated_questions):    
    #                 if st.checkbox(f"Question {i+1}", value=True):
    #                     selected_questions.append(question)
    #                 st.write(question)    
                
    #             # Store generated questions in session state
    #             if 'history' not in st.session_state:
    #                 st.session_state.history = []
    #             st.session_state.history.append({"topic": topic_name, "questions": selected_questions})

    #             st.session_state.generated_questions = selected_questions
    #             st.session_state.topic_name = topic_name
    #             st.success(f"Generated {len(selected_questions)} questions. Please proceed to add them to a question bank.")
    #         else:    
    #             st.error("Failed to generate questions")        

    #     # New feature: View Generated Questions History
    #     st.subheader("View Generated Questions History")
    #     topic_questions = get_generated_questions_history()  # Retrieve the history

    #     if topic_questions:
    #         selected_topic = st.selectbox("Select Topic", options=list(topic_questions .keys()))
    #         if selected_topic:
    #             questions = topic_questions[selected_topic]
    #             st.write(f"Generated Questions for Topic: {selected_topic}")
    #             for i, question in enumerate(questions, 1):
    #                 st.write(f"{i}. {question}")
    #     else:
    #         st.info("No generated questions history available.")

    #     if 'generated_questions' in st.session_state and st.session_state.generated_questions:
    #         st.subheader("Add Questions to Question Bank")
            
    #         existing_question_banks = get_all_question_banks()
    #         qb_options = ["Create New Question Bank"] + [f"ID: {qb['id']} - {qb['technology']} - {qb['difficulty']}" for qb in existing_question_banks]
    #         selected_qb = st.selectbox("Select Question Bank", options=qb_options)
            
    #         if st.button("Add Questions to Selected Bank"):
    #             if selected_qb == "Create New Question Bank":
    #                 new_qb_technology = st.text_input("Enter technology for new question bank")
    #                 new_qb_difficulty = st.selectbox("Select difficulty for new question bank", ["Easy", "Medium", "Hard"])
    #                 if st.button("Create and Add"):
    #                     new_qb_id = create_new_question_bank(new_qb_technology, new_qb_difficulty, st.session_state.generated_questions)
    #                     if new_qb_id:
    #                         st.success(f"Created new question bank with ID: {new_qb_id} and added selected questions.")
    #                         st.session_state.generated_qb_id = new_qb_id  # Store the new question bank ID
    #                     else:
    #                         st.error("Failed to create new question bank.")
    #             else:
    #                 qb_id = int(selected_qb.split('-')[0].split(':')[1].strip())
    #                 if add_questions_to_question_bank(qb_id, st.session_state.topic_name, st.session_state.generated_questions):
    #                     st.success(f"Questions added to question bank ID: {qb_id}")
    #                 else:
    #                     st.error("Failed to add questions to the selected question bank.")

    #             # Clear the generated questions from session state
    #             del st.session_state.generated_questions
    #             del st.session_state.topic_name

    elif selected == "Employee Performance":
        st.subheader("Employee Performance 📈")
        employees = get_all_users()

        if employees:
            selected_employee = st.selectbox(
                "Select Employee",
                options=[employee['username'] for employee in employees],
                key="employee_performance_select"
            )

            if selected_employee:
                # Fetch assessment results for the selected employee
                assessment_results = get_assessment_results(selected_employee)
                if assessment_results:
                    # Prepare data for the table
                    performance_data = []
                    for result in assessment_results:
                        performance_data.append({
                            'Question Bank ID': str(result['question_bank_id']), # Convert ObjectId to string
                            'Score': result['score'],
                            'Completed At': result['completed_at']
                        })

                    # Convert to DataFrame for better visualization
                    performance_df = pd.DataFrame(performance_data)

                    # Display summary metrics
                    st.subheader(f"Summary Statistics for {selected_employee}")
                    total_assessments = len(performance_df)
                    avg_score = performance_df['Score'].mean() if total_assessments > 0 else 0
                    best_score = performance_df['Score'].max() if total_assessments > 0 else 0

                    col1, col2, col3 = st.columns(3)
                    with col1:
                        st.metric("Total Assessments", total_assessments)
                    with col2:
                        st.metric("Average Score", f"{avg_score:.1f}")
                    with col3:
                        st.metric("Best Score", best_score)

                    # Display the performance data in a styled table
                    st.write(f"Performance Data for {selected_employee}:")
                    st.dataframe(performance_df.style.highlight_max(axis=0))  # Highlight max scores

                    # Convert 'Completed At' column to datetime for sorting
                    performance_df['Completed At'] = pd.to_datetime(performance_df['Completed At'])
                    performance_df.sort_values('Completed At', inplace=True)

                    # Create visualizations
                    st.subheader("Performance Over Time")

                    # Line chart for scores over time
                    fig_line = px.line(performance_df, x='Completed At', y='Score', 
                                    title='Score Over Time', markers=True)
                    st.plotly_chart(fig_line)

                    # Bar chart for scores by question bank
                    fig_bar = px.bar(performance_df, x='Question Bank ID', y='Score', 
                                    title='Scores by Question Bank', text='Score')
                    st.plotly_chart(fig_bar)

                    # Convert figures to HTML format for download
                    fig_line_html = fig_line.to_html(full_html=False)
                    fig_bar_html = fig_bar.to_html(full_html=False)

                    # Provide download buttons
                    st.download_button(label="Download Line Chart as HTML", data=fig_line_html, 
                                    file_name=f"{selected_employee}_performance_over_time.html", mime="text/html")

                    st.download_button(label="Download Bar Chart as HTML", data=fig_bar_html, 
                                    file_name=f"{selected_employee}_score_by_question_bank.html", mime="text/html")

                    st.download_button(label="Download Performance Data as CSV", 
                                    data=performance_df.to_csv(index=False), 
                                    file_name=f"{selected_employee}_performance.csv", mime="text/csv")

                else:
                    st.info("No assessment results available for this employee.")
        else:
            st.info("No employees available.")

    # Display content based on the selected option
    if selected == "Generate Questions":
        # Horizontal menu for question generation methods
        question_generation_method = option_menu(
            menu_title=None,  # required
            options=["Generate Questions by Topic", "Generate Questions from Prompt"],  # required
            icons=["book", "pencil"],  # optional
            menu_icon="cast",  # optional
            default_index=0,  # optional
            orientation="horizontal",
        )

        if question_generation_method == "Generate Questions by Topic":
            st.subheader("Generate Questions by Topic 🚀")
            topic_name = st.text_input("Enter Topic Name", key="topic_input")  # Input for topic name
            num_questions = st.number_input("Number of Questions to Generate", min_value=1, value=5, key="num_questions_input")  # Input for number of questions
            
            # Dropdown for selecting question type
            question_type = st.selectbox("Select Question Type", ["Multiple Choice", "Subjective", "Fill in the Blanks"], key="question_type_select")

            if st.button("Generate Questions", key="generate_topic_questions_button"):
                if topic_name:
                    try:
                        prompt = f"Generate {num_questions} {question_type.lower()} questions based on the topic: {topic_name}."
                        generated_questions = model.generate_content(prompt)  # Using the model to generate questions
                        questions_text = generated_questions.text.strip()  # Extracting the generated text

                        # Displaying the generated questions
                        st.write("Generated Questions:")
                        questions_list = questions_text.split('\n')
                        selected_questions = []
                        
                        for i, question in enumerate(questions_list, 1):
                            if question.strip():  # Only show non-empty questions
                                # Display each question with a checkbox
                                if st.checkbox(f"Question {i}: {question.strip()}", value=True):
                                    selected_questions.append(question.strip())

                        # Store generated questions in session state
                        if 'history' not in st.session_state:
                            st.session_state.history = []
                        st.session_state.history.append({"topic": topic_name, "questions": selected_questions})

                        st.session_state.generated_questions = selected_questions
                        st.session_state.topic_name = topic_name
                        st.success(f"Generated {len(selected_questions)} questions. Please proceed to add them to a question bank.")
                    except Exception as e:
                        st.error(f"Error generating questions: {e}")
                else:
                    st.error("Please enter a topic name.")

            # New feature: Add Questions to Question Bank
            if 'generated_questions' in st.session_state and st.session_state.generated_questions:
                st.subheader("Add Questions to Question Bank")
                
                existing_question_banks = get_all_question_banks()
                qb_options = ["Create New Question Bank"] + [f"ID: {str(qb['_id'])} - {qb['technology']} - {qb['difficulty']}" for qb in existing_question_banks]
                selected_qb = st.selectbox("Select Question Bank", options=qb_options)
                
                if st.button("Add Questions to Selected Bank"):
                    if selected_qb == "Create New Question Bank":
                        new_qb_technology = st.text_input("Enter technology for new question bank")
                        new_qb_difficulty = st.selectbox("Select difficulty for new question bank", ["Easy", "Medium", "Hard"])
                        if st.button("Create and Add"):
                            new_qb_id = create_new_question_bank(new_qb_technology, new_qb_difficulty, st.session_state.generated_questions)
                            if new_qb_id:
                                st.success(f"Created new question bank with ID: {new_qb_id} and added selected questions.")
                                st.session_state.generated_qb_id = new_qb_id  # Store the new question bank ID
                            else:
                                st.error("Failed to create new question bank.")
                    else:
                        qb_id = ObjectId(selected_qb.split('-')[0].split(':')[1].strip()) # Convert to ObjectId
                        if add_questions_to_question_bank(qb_id, st.session_state.topic_name, st.session_state.generated_questions):
                            st.success(f"Questions added to question bank ID: {qb_id}")
                        else:
                            st.error("Failed to add questions to the selected question bank.")

                    # Clear the generated questions from session state
                    del st.session_state.generated_questions
                    del st.session_state.topic_name

        elif question_generation_method == "Generate Questions from Prompt":
            st.subheader("Generate Questions from Prompt ✍️")    
            topic_name = st.text_input("Enter Topic Name")    
            prompt = st.text_area("Enter a paragraph to generate questions")    
            question_type = st.selectbox("Select Question Type", ["Multiple Choice", "Subjective", "Fill in the Blanks"])    
            difficulty_level = st.selectbox("Select Difficulty Level", ["Easy", "Medium", "Hard"])  
            num_questions = st.number_input("Number of Questions to Generate", min_value=1, value=10)  
            
            if st.button("Generate Questions"):    
                generated_questions = generate_questions_from_prompt(prompt, question_type, difficulty_level, num_questions, topic_name)    
                if generated_questions:    
                    st.write("Generated Questions:")    
                    selected_questions = []
                    for i, question in enumerate(generated_questions):    
                        if st.checkbox(f"Question {i+1}", value=True):
                            selected_questions.append(question)
                        st.write(question)    
                    
                    # Store generated questions in session state
                    if 'history' not in st.session_state:
                        st.session_state.history = []
                    st.session_state.history.append({"topic": topic_name, "questions": selected_questions})

                    st.session_state.generated_questions = selected_questions
                    st.session_state.topic_name = topic_name
                    st.success(f"Generated {len(selected_questions)} questions. Please proceed to add them to a question bank.")
                else:    
                    st.error("Failed to generate questions")        

            # New feature: View Generated Questions History
            st.subheader("View Generated Questions History")
            topic_questions = get_generated_questions_history()  # Retrieve the history

            if topic_questions:
                selected_topic = st.selectbox("Select Topic", options=list(topic_questions.keys()))
                if selected_topic:
                    questions = topic_questions[selected_topic]
                    st.write(f"Generated Questions for Topic: {selected_topic}")
                    for i, question in enumerate(questions, 1):
                        st.write(f"{i}. {question}")
            else:
                st.info("No generated questions history available.")

            if 'generated_questions' in st.session_state and st.session_state.generated_questions:
                st.subheader("Add Questions to Question Bank")
                
                existing_question_banks = get_all_question_banks()
                qb_options = ["Create New Question Bank"] + [f"ID: {str(qb['_id'])} - {qb['technology']} - {qb['difficulty']}" for qb in existing_question_banks]
                selected_qb = st.selectbox("Select Question Bank", options=qb_options)
                
                if st.button("Add Questions to Selected Bank"):
                    if selected_qb == "Create New Question Bank":
                        new_qb_technology = st.text_input("Enter technology for new question bank")
                        new_qb_difficulty = st.selectbox("Select difficulty for new question bank", ["Easy", "Medium", "Hard"])
                        if st.button("Create and Add"):
                            new_qb_id = create_new_question_bank(new_qb_technology, new_qb_difficulty, st.session_state.generated_questions)
                            if new_qb_id:
                                st.success(f"Created new question bank with ID: {new_qb_id} and added selected questions.")
                                st.session_state.generated_qb_id = new_qb_id  # Store the new question bank ID
                            else:
                                st.error("Failed to create new question bank.")
                    else:
                        qb_id = ObjectId(selected_qb.split('-')[0].split(':')[1].strip()) # Convert to ObjectId
                        if add_questions_to_question_bank(qb_id, st.session_state.topic_name, st.session_state.generated_questions):
                            st.success(f"Questions added to question bank ID: {qb_id}")
                        else:
                            st.error("Failed to add questions to the selected question bank.")

        # Clear the generated questions from session state
                        del st.session_state.generated_questions
                        del st.session_state.topic_name
        # Display content based on the selected option
    
    elif selected == "Chatbot":
        # Display chatbot interface at the top
        st.subheader("Chat with the AI Trainer 🤖")
        
        # Initialize the message state if not exists
        if "msg" not in st.session_state:
            st.session_state.msg = ""
        
        # Create a container for the chat interface
        chat_container = st.container()

        # Define avatars
        user_avatar = "https://static.vecteezy.com/system/resources/previews/009/664/418/non_2x/people-user-team-transparent-free-png.png"
        ai_avatar = "https://thumbs.dreamstime.com/b/chatbot-logo-messenger-ai-robot-icon-vector-illustration-277900892.jpg"

        def clear_text():
            st.session_state.msg = st.session_state.user_input
            st.session_state.user_input = ""

        with chat_container:
            # Display chat messages
            for chat in st.session_state.chat_history:
                if chat['role'] == 'assistant':
                    # Chatbot message with avatar
                    st.markdown(
                        f"<div style='display: flex; align-items: center; margin: 5px 0;color:black'>"
                        f"<img src='{ai_avatar}' style='width: 40px; height: 40px; border-radius: 50%; margin-right: 10px;'>"
                        f"<div style='background-color: #e1ffc7; padding: 10px; border-radius: 10px; max-width: 80%;'>"
                        f"<strong>AI:</strong> {chat['content']}</div></div>",
                        unsafe_allow_html=True
                    )
                else:
                    # User message with avatar
                    st.markdown(
                        f"<div style='display: flex; align-items: center; margin: 5px 0; justify-content: flex-end;color:black'>"
                        f"<div style='background-color: #dcf8c6; padding: 10px; border-radius: 10px; max-width: 80%; margin-left: auto;'>"
                        f"<strong>You:</strong> {chat['content']}</div>"
                        f"<img src='{user_avatar}' style='width: 40px; height: 40px; border-radius: 50%; margin-left: 10px;'>"
                        f"</div>", 
                        unsafe_allow_html=True
                    )

            # Input field for user to enter a prompt
            st.text_input("Type your message here...", key="user_input", placeholder="Type a message...", on_change=clear_text)

            if st.session_state.msg:  # Only process if there's a message
                # Append user input to chat history
                st.session_state.chat_history.append({"role": "user", "content": st.session_state.msg})

                try:
                    # Generate AI response
                    prompt = f"You are an AI assistant for trainers. Respond to the following message: {st.session_state.msg}"
                    response = model.generate_content(prompt)
                    
                    # Handle the response properly for Gemini model
                    if hasattr(response, 'parts'):
                        ai_response = ''.join(part.text for part in response.parts)
                    else:
                        ai_response = response.candidates[0].content.parts[0].text
                    
                    # Append AI response to chat history
                    st.session_state.chat_history.append({"role": "assistant", "content": ai_response})
                except Exception as e:
                    st.error(f"Error generating response: {str(e)}")
                    ai_response = "I apologize, but I encountered an error. Please try again."
                    st.session_state.chat_history.append({"role": "assistant", "content": ai_response})
                
                # Clear the message state
                st.session_state.msg = ""
                
                # Rerun the app to display the new messages
                st.rerun()

                    
    
    

        
    notifications = get_notifications("trainer", None)  # Get notifications for trainer
    display_notifications(notifications, "trainer")  # Display notifications in the sidebar

    if notifications:  
        st.sidebar.write("Notifications:")  
        for notification in notifications:  
            st.sidebar.write(notification['message'])  
    else:  
        st.sidebar.write("No notifications available.")



    
# Download NLTK data
nltk.download('vader_lexicon')

# Load .env file (if used locally for development)
# from dotenv import load_dotenv
# load_dotenv()

# Set Google API Key from environment variable
os.environ["GOOGLE_API_KEY"] = os.getenv("GOOGLE_API_KEY", "AIzaSyA9F4eR8dUfZQdwUcL-le9RP-G3Mssn7T8") # Fallback for local testing

# Configure Google Generative AI model
genai.configure(api_key=os.getenv('GOOGLE_API_KEY'))
model = genai.GenerativeModel('gemini-2.0-flash')


# Modified save_question_bank to use MongoDB
def save_question_bank(technology, topics, questions, difficulty, correct_answers, question_type, options=None):
    db = create_connection()
    if db is None:
        return None

    try:
        # Prepare document for question_banks collection
        qb_doc = {
            "technology": technology,
            "topics": topics, # topics is already a list of strings
            "questions": questions, # questions is already a single string
            "difficulty": difficulty,
            "question_type": question_type,
            "options": options, # options is already a single string
            "created_at": datetime.now()
        }
        
        # Insert into question_banks collection
        result_qb = db.question_banks.insert_one(qb_doc)
        question_bank_id = result_qb.inserted_id # MongoDB's _id

        # Prepare document for question_answers collection
        answer_doc = {
            "question_bank_id": question_bank_id, # Link to the question bank
            "answer_data": correct_answers # correct_answers is already a single string
        }
        
        # Insert into question_answers collection
        db.question_answers.insert_one(answer_doc)
        
        return str(question_bank_id) # Return as string for consistency with app logic
    except OperationFailure as e:
        print(f"MongoDB operation error: {e}")
        return None
    except Exception as e:
        print(f"General error in save_question_bank: {e}")
        return None

# MongoDB connection
def create_connection():
    try:
        # Replace with your MongoDB connection string from environment variable or fallback
        mongo_uri = os.getenv("MONGO_URI", "mongodb://localhost:27017/")
        client = MongoClient(mongo_uri) 
        db = client["final_mongodb"] # Your database name
        return db
    except ConnectionFailure as e:
        st.error(f"Error connecting to MongoDB: {e}")
        return None


def ascii_to_string(ascii_list):
    """Convert a list of ASCII values to a string."""
    return ''.join(chr(num) for num in ascii_list)

    # Example usage when retrieving topics
    curricula = get_all_curricula()
    for curriculum in curricula:
        # Assuming 'topics' is a list of ASCII values in the curricula
        if isinstance(curriculum['topics'], list):  # Check if topics is a list of ASCII values
            curriculum['topics'] = ascii_to_string(curriculum['topics'])  # Convert to string


def format_timestamp(timestamp_str):
    """
    Format timestamp for display
    """
    if timestamp_str == 'N/A':
        return 'N/A'
    try:
        # Parse the timestamp string to datetime
        timestamp = datetime.datetime.strptime(timestamp_str, "%Y-%m-%d %H:%M:%S")
        # Format it to a more readable format
        return timestamp.strftime("%b %d, %Y %I:%M %p")
    except (ValueError, TypeError):
        return 'N/A'



# Utility functions
def extract_text_from_file(file):
    file_extension = os.path.splitext(file.name)[1].lower()
    text = ""

    try:
        if file_extension == '.pdf':
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        elif file_extension == '.docx':
            doc = docx.Document(file)
            text = "\n".join([para.text for para in doc.paragraphs])
        elif file_extension == '.txt':
            text = file.getvalue().decode('utf-8')
        elif file_extension in ['.ppt', '.pptx']:
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text += shape.text + "\n"
        elif file_extension == '.csv':
            csv_data = pd.read_csv(file)
            text = csv_data.to_string(index=False)
        else:
            # For other file types, attempt to read as text
            try:
                text = file.getvalue().decode('utf-8')
            except UnicodeDecodeError:
                raise ValueError(f"Unable to extract text from {file_extension} file.")
    except Exception as e:
        raise ValueError(f"Error processing {file_extension} file: {str(e)}")

    # Clean the extracted text
    cleaned_text = clean_text(text)
    return cleaned_text

def clean_text(text):
    # Remove non-printable characters and control characters
    text = re.sub(r'[^\x20-\x7E]+', ' ', text)  # Keep only printable ASCII characters
    
    # Normalize whitespace
    text = re.sub(r'\s+', ' ', text)  # Replace multiple spaces with a single space
    text = text.strip()  # Trim leading and trailing whitespace

    return text


def create_new_question_bank(technology, difficulty, questions):
    db = create_connection()
    if db is None:
        return None

    try:
        questions_text = '\n'.join(questions)
        new_qb_doc = {
            "technology": technology,
            "difficulty": difficulty,
            "questions": questions_text,
            "created_at": datetime.now()
        }
        result = db.question_banks.insert_one(new_qb_doc)
        return str(result.inserted_id) # Return the ObjectId as a string
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None

def generate_questions(text, num_questions=5, question_type="multiple-choice"):
    if question_type == "multiple-choice":
        prompt = f"Generate {num_questions} multiple-choice questions based on the following text:\n\n{text}\n\nProvide the questions and options in the following format:\n\nQ1: [Question]\nA) [Option 1]\nB) [Option 2]\nC) [Option 3]\nD) [Option 4]\n\nQ2: [Question]\nA) [Option 1]\nB) [Option 2]\nC) [Option 3]\nD) [Option 4]\n\n..."
    elif question_type == "subjective":
        prompt = f"Generate {num_questions} subjective questions based on the following text:\n\n{text}\n\nProvide the questions in the following format:\n\nQ1: [Question]\n\nQ2: [Question]\n\n..."
    elif question_type == "fill-in-the-blank":
        prompt = f"Generate {num_questions} fill-in-the-blank questions based on the following text:\n\n{text}\n\nProvide the questions and correct answers in the following format:\n\nQ1: [Question]\nA: [Correct Answer]\n\nQ2: [Question]\nA: [Correct Answer]\n\n..."
    else:
        raise ValueError("Invalid question type")

    response = model.generate_content(prompt)
    generated_text = response.text

    questions = []
    options = []
    correct_answers = []

    lines = [line.strip() for line in generated_text.split('\n') if line.strip()]

    i = 0
    while i < len(lines):
        if lines[i].startswith('Q'):
            question = lines[i].split(': ', 1)[1]
            questions.append(question)
            if question_type == "multiple-choice":
                options_list = []
                correct_answer = None
                # Look for options immediately following the question
                for j in range(i + 1, len(lines)):
                    if lines[j].startswith(('A)', 'B)', 'C)', 'D)')):
                        option = lines[j].split(') ', 1)[1]
                        options_list.append(option)
                        if lines[j].startswith('A)'): # Assuming A is always the correct answer for simplicity in parsing
                            correct_answer = option
                    else:
                        break # Stop if a line doesn't start with an option letter
                options.append(options_list)
                correct_answers.append(correct_answer)
                i = j # Move index to the line after the last option processed
            elif question_type == "fill-in-the-blank":
                if i+1 < len(lines) and lines[i+1].startswith('A:'):
                    options.append([lines[i+1].split(': ', 1)[1]])
                    correct_answers.append(lines[i+1].split(': ', 1)[1])
                    i += 2
                else:
                    options.append([""])
                    correct_answers.append("")
                    i += 1
            else: # subjective
                options.append([])
                correct_answers.append("")
                i += 1
        else:
            i += 1

    return questions[:num_questions], options[:num_questions], correct_answers[:num_questions]

# Removed ensure_table_exists as MongoDB handles collection creation implicitly

def review_feedback():
    # Fetch feedback data from MongoDB
    db = create_connection()
    if db is None:
        return []
    
    try:
        feedback_data = list(db.feedback.find({}))
        # Convert ObjectId to string for compatibility with DataFrame
        for item in feedback_data:
            if '_id' in item:
                item['id'] = str(item['_id'])
                del item['_id'] # Remove ObjectId if not needed for display
            if 'question_bank_id' in item and isinstance(item['question_bank_id'], ObjectId):
                item['question_bank_id'] = str(item['question_bank_id'])
        return feedback_data
    except OperationFailure as e:
        st.error(f"Error retrieving feedback: {e}")
        return []

def analyze_sentiment(text):
    sia = SentimentIntensityAnalyzer()
    sentiment_score = sia.polarity_scores(text)['compound']
    if sentiment_score > 0.05:
        return 'Positive'
    elif sentiment_score < -0.05:
        return 'Negative'
    else:
        return 'Neutral'

def display_questions(questions, options, correct_answers):
    for i, question in enumerate(questions):
        st.write(question)
        if options[i]:
            st.write("Options:")
            for j, option in enumerate(options[i]):
                st.write(f"{chr(65+j)}) {option}")
        st.write(f"Correct Answer: {correct_answers[i]}")
        st.write("")
# User Authentication Functions
def login_user(username, password):
    db = create_connection()
    if db is None:
        return None

    user = db.users.find_one({"username": username})

    if user and check_password_hash(user['password'], password):
        # Convert ObjectId to string for session state
        user['_id'] = str(user['_id'])
        return user
    return None

def register_user(email, username, password, role):
    # Validate email format
    email_regex = r'^[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+$'
    if not re.match(email_regex, email):
        st.error("Invalid email format. Please enter a valid email address.")
        return False

    db = create_connection()
    if db is None:
        return False

    existing_user = db.users.find_one({"username": username})

    if existing_user:
        return False

    hashed_password = generate_password_hash(password)
    user_data = {
        "email": email,
        "username": username,
        "password": hashed_password,
        "role": role
    }
    try:
        db.users.insert_one(user_data)
        return True
    except OperationFailure as e:
        st.error(f"Database error during registration: {e}")
        return False
# Administrator Functions
def get_system_stats():
    db = create_connection()
    if db is None:
        return None

    stats = {}
    try:
        # Get count of documents in each collection
        stats['users'] = db.users.count_documents({})
        stats['question_banks'] = db.question_banks.count_documents({})
        stats['learning_plans'] = db.learning_plans.count_documents({})
        stats['feedback'] = db.feedback.count_documents({})
        return stats
    except OperationFailure as e:
        st.error(f"Error retrieving system stats: {e}")
        return None

def get_all_users():
    db = create_connection()
    if db is None:
        return []

    try:
        users_cursor = db.users.find({}, {"username": 1, "email": 1, "role": 1, "_id": 0})
        users = list(users_cursor)
        return users
    except OperationFailure as e:
        st.error(f"Error retrieving all users: {e}")
        return []

def update_user_role(username, new_role):
    db = create_connection()
    if db is None:
        return False

    try:
        user_exists = db.users.find_one({"username": username})
        if not user_exists:
            st.error("User does not exist.")
            return False

        result = db.users.update_one({"username": username}, {"$set": {"role": new_role}})
        return result.modified_count > 0
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False

# Trainer Functions
def upload_curriculum(technology, topics, content):   
  db = create_connection()   
  if db is None:   
    return False   
   
  try:   
    topics_str = ','.join(topics)   
    
    # Determine the content type (file-like object or string)   
    if hasattr(content, 'read'):   
      # File-like object   
      content_text = content.read().decode('utf-8')   
    else:   
      # String   
      content_text = content   
   
    # Generate questions from the content   
    questions, options, correct_answers = generate_questions(content_text)   
   
    # Convert questions, options and correct_answers to strings   
    questions_str = '|||'.join(questions)   
    options_str = '|||'.join(['###'.join(option) for option in options]) # Using '###' as separator for options   
    correct_answers_str = '|||'.join([','.join(map(str, ans)) if isinstance(ans, list) else str(ans) for ans in correct_answers])   
   
    # Insert or update the curriculum using upsert
    curriculum_doc = {
        "technology": technology,
        "topics": topics_str,
        "filename": "curriculum_" + technology + ".txt",
        "content": content_text
    }
    db.curriculum.update_one({"technology": technology}, {"$set": curriculum_doc}, upsert=True)
   
    # Insert the generated questions into the generated_question_files collection
    generated_questions_doc = {
        "technology": technology,
        "topics": topics_str,
        "questions": questions_str,
        "options": options_str,
        "correct_answers": correct_answers_str,
        "created_at": datetime.now()
    }
    db.generated_question_files.insert_one(generated_questions_doc)

    return True   
  except OperationFailure as err:   
    st.error(f"Database error: {err}")   
    return False   
  except Exception as e:
    st.error(f"Error in upload_curriculum: {e}")
    return False


def get_curriculum_text(technology):
    db = create_connection()
    if db is None:
        return None

    try:
        result = db.curriculum.find_one({"technology": technology}, {"topics": 1, "_id": 0})
        if result:
            return result.get('topics')
        else:
            st.error(f"No curriculum content found for technology: {technology}")
            return None
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None

def save_question_bank(technology, topics, questions, difficulty, correct_answers, question_type, options=None):  # ✅ Correct one

    db = create_connection()
    if db is None:
        return False

    try:
        # Prepare document for question_banks collection
        qb_doc = {
            "technology": technology,
            "topics": topics, # topics is already a string
            "questions": questions, # questions is already a single string
            "difficulty": difficulty,
            "question_type": question_type,
            "options": options, # options is already a single string
            "created_at": datetime.now()
        }
        
        # Insert into question_banks collection
        result_qb = db.question_banks.insert_one(qb_doc)
        question_bank_id = result_qb.inserted_id # MongoDB's _id

        # Prepare document for question_answers collection
        answer_doc = {
            "question_bank_id": question_bank_id, # Link to the question bank
            "answer_data": correct_answers # correct_answers is already a single string
        }
        
        # Insert into question_answers collection
        db.question_answers.insert_one(answer_doc)
        
        return str(question_bank_id) # Return as string for consistency with app logic
        
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False
    except Exception as e:
        st.error(f"General error in save_question_bank: {e}")
        return False

def get_topics_for_technology(technology):
    db = create_connection()
    if db is None:
        return None

    try:
        result = db.curriculum.find_one({"technology": technology}, {"topics": 1, "_id": 0})
        if result and 'topics' in result:
            return result['topics'].split(',')
        else:
            st.error(f"No topics found for technology: {technology}")
            return None
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None

def get_all_curricula():
    db = create_connection()
    if db is None:
        return None

    try:
        curricula_cursor = db.curriculum.find({}, {"technology": 1, "topics": 1}) # Include _id for potential future use
        curricula = []
        for doc in curricula_cursor:
            doc['id'] = str(doc['_id']) # Add string version of _id
            curricula.append(doc)
        return curricula
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None

def get_all_question_banks():
    db = create_connection()
    if db is None:
        return None

    try:
        # Fetch all documents from the question_banks collection
        question_banks_cursor = db.question_banks.find({})
        question_banks = list(question_banks_cursor)
        return question_banks
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None

# Employee Functions
def get_learning_plan(username):
    db = create_connection()
    if db is None:
        return None

    learning_plan = db.learning_plans.find_one({"username": username})
    return learning_plan

def submit_feedback(username, question_bank_id, feedback_text, rating, feedback_type):
    sentiment = analyze_sentiment(feedback_text)  # Optional: Analyze sentiment
    db = create_connection()
    if db is None:
        return False

    try:
        feedback_doc = {
            "username": username,
            "question_bank_id": ObjectId(question_bank_id) if question_bank_id else None, # Store as ObjectId
            "feedback_text": feedback_text,
            "rating": rating,
            "sentiment": sentiment,
            "feedback_type": feedback_type,
            "created_at": datetime.now()
        }
        db.feedback.insert_one(feedback_doc)

        # Prepare notification message
        feedback_summary = f"New feedback received from {username}. Type: {feedback_type}. Rating: {rating}. Feedback: {feedback_text}"

        # Send notifications based on feedback type
        if feedback_type in ["User  Experience"] and rating >= 3:
            send_notification("admin", feedback_summary, username)
        elif feedback_type in ["Question Bank Feedback", "Assessment Feedback"]:
            send_notification("admin", feedback_summary, username)
            send_notification("trainer", feedback_summary, username)

        return True
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False


def take_assessment():
    st.subheader("Take Assessment")
    question_banks = get_all_question_banks()
    if not question_banks:
        st.info("No question banks available yet.")
    else:
        selected_qb = st.selectbox(
            "Select Question Bank",
            options=[(str(qb['_id']), f"{qb['technology']} - {qb['difficulty']}") for qb in question_banks],
            format_func=lambda x: f"ID: {x[0]} - {x[1]}",
            key="take_assessment_qb_select"
        )

        if selected_qb:
            qb_id_str, _ = selected_qb
            qb_id = ObjectId(qb_id_str) # Convert to ObjectId
            qb_details = next((qb for qb in question_banks if qb['_id'] == qb_id), None)
            if qb_details:
                questions = qb_details.get('questions', '').split('\n')
                options = qb_details.get('options', '').split('###') if qb_details.get('options') else []
                correct_answers_str = get_correct_answers(qb_id) # This returns a list of strings
                
                question_type = qb_details.get('question_type')

                score = 0
                user_answers = []
                for i, question in enumerate(questions):
                    if not question.strip(): # Skip empty questions
                        continue

                    st.write(f"**Q{i+1}:** {question.strip()}")
                    
                    if question_type == "multiple-choice" and options:
                        # Assuming options are stored as 'Opt1###Opt2###Opt3###Opt4' per question
                        # Need to parse options for the current question correctly
                        # For now, let's assume options are structured to match questions
                        if i < len(options): # Ensure we have options for this question
                            current_options = options[i].split('###')
                            if current_options:
                                answer = st.radio("Select an option", current_options, key=f"question_{i}")
                                user_answers.append(answer)
                            else:
                                user_answers.append("") # No options, no answer
                        else:
                            user_answers.append("") # No options for this question
                    elif question_type == "fill-in-the-blank":
                        answer = st.text_input("Enter your answer", key=f"question_{i}")
                        user_answers.append(answer)
                    elif question_type == "subjective":
                        answer = st.text_area("Enter your answer", key=f"question_{i}")
                        user_answers.append(answer)
                    else:
                        user_answers.append("") # For cases where question_type is not set or options are missing

                if st.button("Submit"):
                    total_questions_answered = 0
                    correct_answers_count = 0
                    for i, user_answer in enumerate(user_answers):
                        if i < len(correct_answers_str) and user_answer.strip().lower() == correct_answers_str[i].strip().lower():
                            st.success(f"Q{i+1}: Correct!")
                            correct_answers_count += 1
                        elif i < len(correct_answers_str):
                            st.error(f"Q{i+1}: Incorrect. Correct answer: {correct_answers_str[i].strip()}")
                        total_questions_answered += 1

                    st.write(f"Your score is {correct_answers_count} out of {total_questions_answered}")

                    # Save the assessment result
                    save_assessment_result(st.session_state.user['username'], qb_id, correct_answers_count)

def get_available_question_banks(username):
    db = create_connection()
    if db is None:
        return []

    try:
        # Get technologies from learning plans for the user
        learning_plans_cursor = db.learning_plans.find({"username": username}, {"technology": 1, "_id": 0})
        technologies = [lp['technology'] for lp in learning_plans_cursor if 'technology' in lp]

        # Get IDs of question banks already completed by the user
        completed_assessments_cursor = db.assessments.find({"username": username}, {"question_bank_id": 1, "_id": 0})
        completed_qb_ids = [a['question_bank_id'] for a in completed_assessments_cursor if 'question_bank_id' in a]

        # Find question banks matching learning plan technologies and not yet completed
        query = {
            "technology": {"$in": technologies},
            "_id": {"$nin": completed_qb_ids}
        }
        question_banks_cursor = db.question_banks.find(query, {"technology": 1, "topics": 1})
        
        # Convert ObjectId to string for 'id' field in the returned dictionary
        question_banks = []
        for qb in question_banks_cursor:
            qb['id'] = str(qb['_id'])
            question_banks.append(qb)
        
        return question_banks
    except OperationFailure as e:
        st.error(f"Error retrieving available question banks: {e}")
        return []

def get_completed_assessments(username):
    db = create_connection()
    if db is None:
        return []

    try:
        # Fetch assessments for the user
        assessments_cursor = db.assessments.find({"username": username}).sort("completed_at", -1)
        completed_assessments = []

        for assessment in assessments_cursor:
            # Fetch corresponding question bank details
            qb_details = db.question_banks.find_one({"_id": assessment['question_bank_id']})
            
            if qb_details:
                total_questions = len(qb_details.get('questions', '').split('\n')) if qb_details.get('questions') else 0
                
                percentage = round((assessment['score'] / total_questions) * 100, 2) if total_questions > 0 else 0.0

                completed_assessments.append({
                    'id': str(assessment['_id']), # Convert ObjectId to string
                    'question_bank_id': str(assessment['question_bank_id']), # Convert ObjectId to string
                    'technology': qb_details.get('technology'),
                    'difficulty': qb_details.get('difficulty'),
                    'score': assessment.get('score'),
                    'question_type': qb_details.get('question_type'),
                    'completed_at': assessment.get('completed_at').strftime('%Y-%m-%d %H:%M:%S') if assessment.get('completed_at') else 'N/A',
                    'total_questions': total_questions,
                    'percentage': percentage
                })
        return completed_assessments
    except OperationFailure as e:
        st.error(f"Error retrieving completed assessments: {e}")
        return []


def admin_dashboard():  
      
    with st.sidebar:
      # Create a sidebar for navigation using option_menu
      selected_tab = option_menu(
          menu_title="Admin Dashboard",  # required
          options=["System Stats", 
                  "User  Management", 
                  "Reports", 
                  "Employee Performance"],  # required
          icons=["bar-chart", "people", "file-earmark-text", "person-check"],  # optional
          menu_icon="cast",  # optional
          default_index=0,  # optional
          orientation="vertical",
      )

    # Display the selected tab content
    if selected_tab == "System Stats":
        st.subheader("System Statistics 📊")  
        stats = get_system_stats()  
        if stats:  
            st.write(f"Total Users: {stats.get('users', 0)}")  
            st.write(f"Total Question Banks: {stats.get('question_banks', 0)}")  
            st.write(f"Total Learning Plans: {stats.get('learning_plans', 0)}")  
            st.write(f"Total Feedback Entries: {stats.get('feedback', 0)}")  

            # Add system details  
            st.subheader("System Details 🖥️")  
            st.write(f"Operating System: {platform.system()}")  
            st.write(f"Platform: {platform.platform()}")  
            st.write(f"Processor: {platform.processor()}")  
            st.write(f"Python Version: {platform.python_version()}")  
        else:  
            st.error("Failed to retrieve system statistics") 

    elif selected_tab == "User  Management":
        st.subheader("User  Management 👤")
        
        # Check if users are already loaded in session state
        if 'users' not in st.session_state:
            st.session_state.users = get_all_users()  # Load users for the first time

        # Create a layout for the refresh button
        col1, col2 = st.columns([4, 1])  # Adjust column widths

        with col1:
            st.write("")  # Empty space for alignment

        with col2:
            # Refresh button with an icon
            if st.button("🔄", key="refresh_users", help="Refresh User List"):
                st.session_state.users = get_all_users()  # Refresh the user list

        users = st.session_state.users  # Use the loaded users from session state
        
        if users:
            # Display the user table
            user_table = []
            for user in users:
                user_table.append({
                    'Username': user.get('username'),
                    'Email': user.get('email'),
                    'Role': user.get('role')
                })
            
            user_df = pd.DataFrame(user_table)
            st.table(user_df)

            # CSS to inject for compact layout
            st.markdown("""
                <style>
                .stSelectbox {
                    margin-bottom: 0px;
                }
                .stButton {
                    display: inline-block;
                    margin-right: 10px;
                }
                .user-row {
                    margin-bottom: 10px;
                }
                </style>
            """, unsafe_allow_html=True)

            # Search feature
            st.subheader("Search User")
            search_username = st.text_input("Enter username to search:", "")
            
            if search_username:
                filtered_users = [user for user in users if search_username.lower() in user.get('username', '').lower()]
            else:
                filtered_users = users

            # Display users (filtered or all)
            if filtered_users:
                for user in filtered_users:
                    with st.container():
                        cols = st.columns([2, 3, 3])  # Adjusted column widths
                        
                        with cols[0]:  # Username column
                            st.write(user.get('username'))
                        
                        with cols[1]:  # Role selection column
                            new_role = st.selectbox(
                                f"New Role for {user.get('username')}", 
                                ["None", "Administrator", "Trainer", "Employee"], 
                                key=f"new_role_{user.get('username')}",
                                label_visibility="collapsed"  # Hides the label
                            )
                        
                        with cols[2]:  # Buttons column
                            c1, c2 = st.columns([1, 1])  # Split the column for buttons
                            with c1:
                                if st.button("Update Role", key=f"update_role_button_{user.get('username')}", use_container_width=True):
                                    if update_user_role(user.get('username'), new_role):
                                        st.success(f"Role updated for {user.get('username')}")
                                        # Refresh the user list after update
                                        st.session_state.users = get_all_users()
                                    else:
                                        st.error("Failed to update role")
                            
                            with c2:
                                if st.button("Remove User", key=f"remove_user_button_{user.get('username')}", use_container_width=True):
                                    if remove_user(user.get('username')):
                                        st.success(f"User {user.get('username')} removed successfully")
                                        # Refresh the user list after removal
                                        st.session_state.users = get_all_users()
                                    else:
                                        st.error("Failed to remove user")
            else:
                st.info("No users found with that username.")
        else:
            st.info("No users available.")

    elif selected_tab == "Reports":
        st.subheader("Generate Reports 🔍")
        report_type = st.selectbox("Select Report Type", 
                                    ["User Activity", "Question Bank Usage", 
                                    "Feedback Summary", "Sentiment Analysis", 
                                    "Employee Performance"])

        if st.button("Generate Report"):
            if report_type == "User Activity":
                user_activity_report()
            elif report_type == "Question Bank Usage":
                question_bank_usage_report()
            elif report_type == "Feedback Summary":
                feedback_summary_report()
            elif report_type == "Sentiment Analysis":
                sentiment_analysis_report()
            elif report_type == "Employee Performance":
                employee_performance_report()



    elif selected_tab == "Employee Performance":
        st.subheader("Employee Performance 🎯")
        employees = get_all_users()

        if employees:
            selected_employee = st.selectbox(
                "Select Employee",
                options=[employee['username'] for employee in employees],
                key="employee_performance_select"
            )

            if selected_employee:
                # Fetch assessment results for the selected employee
                assessment_results = get_assessment_results(selected_employee)
                if assessment_results:
                    # Prepare data for the table
                    performance_data = []
                    for result in assessment_results:
                        performance_data.append({
                            'Question Bank ID': str(result['question_bank_id']), # Convert ObjectId to string
                            'Score': result['score'],
                            'Completed At': result['completed_at']
                        })

                    # Convert to DataFrame for better visualization
                    performance_df = pd.DataFrame(performance_data)

                    # Display summary metrics
                    st.subheader(f"Summary Statistics for {selected_employee}")
                    total_assessments = len(performance_df)
                    avg_score = performance_df['Score'].mean() if total_assessments > 0 else 0
                    best_score = performance_df['Score'].max() if total_assessments > 0 else 0

                    col1, col2, col3 = st.columns(3)
                    with col1:
                        st.metric("Total Assessments", total_assessments)
                    with col2:
                        st.metric("Average Score", f"{avg_score:.1f}")
                    with col3:
                        st.metric("Best Score", best_score)

                    # Display the performance data in a styled table
                    st.write(f"Performance Data for {selected_employee}:")
                    st.dataframe(performance_df.style.highlight_max(axis=0))  # Highlight max scores

                    # Convert 'Completed At' column to datetime for sorting
                    performance_df['Completed At'] = pd.to_datetime(performance_df['Completed At'])
                    performance_df.sort_values('Completed At', inplace=True)

                    # Create visualizations
                    st.subheader("Performance Over Time")

                    # Line chart for scores over time
                    fig_line = px.line(performance_df, x='Completed At', y='Score', 
                                    title='Score Over Time', markers=True)
                    st.plotly_chart(fig_line)

                    # Bar chart for scores by question bank
                    fig_bar = px.bar(performance_df, x='Question Bank ID', y='Score', 
                                    title='Scores by Question Bank', text='Score')
                    st.plotly_chart(fig_bar)

                    # Convert figures to HTML format for download
                    fig_line_html = fig_line.to_html(full_html=False)
                    fig_bar_html = fig_bar.to_html(full_html=False)

                    # Provide download buttons
                    st.download_button(label="Download Line Chart as HTML", data=fig_line_html, 
                                    file_name=f"{selected_employee}_performance_over_time.html", mime="text/html")

                    st.download_button(label="Download Bar Chart as HTML", data=fig_bar_html, 
                                    file_name=f"{selected_employee}_score_by_question_bank.html", mime="text/html")

                    st.download_button(label="Download Performance Data as CSV", 
                                    data=performance_df.to_csv(index=False), 
                                    file_name=f"{selected_employee}_performance.csv", mime="text/csv")

                else:
                    st.info("No assessment results available for this employee.")
        else:
            st.info("No employees available.")

    # Display content based on the selected option
        if selected_tab == "Generate Questions":
            # Horizontal menu for question generation methods
            question_generation_method = option_menu(
                menu_title=None,  # required
                options=["Generate Questions by Topic", "Generate Questions from Prompt"],  # required
                icons=["book", "pencil"],  # optional
                menu_icon="cast",  # optional
                default_index=0,  # optional
                orientation="horizontal",
            )

        if question_generation_method == "Generate Questions by Topic":
            st.subheader("Generate Questions by Topic 🚀")
            topic_name = st.text_input("Enter Topic Name", key="topic_input")  # Input for topic name
            num_questions = st.number_input("Number of Questions to Generate", min_value=1, value=5, key="num_questions_input")  # Input for number of questions
            
            # Dropdown for selecting question type
            question_type = st.selectbox("Select Question Type", ["Multiple Choice", "Subjective", "Fill in the Blanks"], key="question_type_select")

            if st.button("Generate Questions", key="generate_topic_questions_button"):
                if topic_name:
                    try:
                        prompt = f"Generate {num_questions} {question_type.lower()} questions based on the topic: {topic_name}."
                        generated_questions = model.generate_content(prompt)  # Using the model to generate questions
                        questions_text = generated_questions.text.strip()  # Extracting the generated text

                        # Displaying the generated questions
                        st.write("Generated Questions:")
                        questions_list = questions_text.split('\n')
                        selected_questions = []
                        
                        for i, question in enumerate(questions_list, 1):
                            if question.strip():  # Only show non-empty questions
                                # Display each question with a checkbox
                                if st.checkbox(f"Question {i}: {question.strip()}", value=True):
                                    selected_questions.append(question.strip())

                        # Store generated questions in session state
                        if 'history' not in st.session_state:
                            st.session_state.history = []
                        st.session_state.history.append({"topic": topic_name, "questions": selected_questions})

                        st.session_state.generated_questions = selected_questions
                        st.session_state.topic_name = topic_name
                        st.success(f"Generated {len(selected_questions)} questions. Please proceed to add them to a question bank.")
                    except Exception as e:
                        st.error(f"Error generating questions: {e}")
                else:
                    st.error("Please enter a topic name.")

            # New feature: Add Questions to Question Bank
            if 'generated_questions' in st.session_state and st.session_state.generated_questions:
                st.subheader("Add Questions to Question Bank")
                
                existing_question_banks = get_all_question_banks()
                qb_options = ["Create New Question Bank"] + [f"ID: {str(qb['_id'])} - {qb['technology']} - {qb['difficulty']}" for qb in existing_question_banks]
                selected_qb = st.selectbox("Select Question Bank", options=qb_options)
                
                if st.button("Add Questions to Selected Bank"):
                    if selected_qb == "Create New Question Bank":
                        new_qb_technology = st.text_input("Enter technology for new question bank")
                        new_qb_difficulty = st.selectbox("Select difficulty for new question bank", ["Easy", "Medium", "Hard"])
                        if st.button("Create and Add"):
                            new_qb_id = create_new_question_bank(new_qb_technology, new_qb_difficulty, st.session_state.generated_questions)
                            if new_qb_id:
                                st.success(f"Created new question bank with ID: {new_qb_id} and added selected questions.")
                                st.session_state.generated_qb_id = new_qb_id  # Store the new question bank ID
                            else:
                                st.error("Failed to create new question bank.")
                    else:
                        qb_id = ObjectId(selected_qb.split('-')[0].split(':')[1].strip()) # Convert to ObjectId
                        if add_questions_to_question_bank(qb_id, st.session_state.topic_name, st.session_state.generated_questions):
                            st.success(f"Questions added to question bank ID: {qb_id}")
                        else:
                            st.error("Failed to add questions to the selected question bank.")

                    # Clear the generated questions from session state
                    del st.session_state.generated_questions
                    del st.session_state.topic_name

        elif question_generation_method == "Generate Questions from Prompt":
            st.subheader("Generate Questions from Prompt ✍️")    
            topic_name = st.text_input("Enter Topic Name")    
            prompt = st.text_area("Enter a paragraph to generate questions")    
            question_type = st.selectbox("Select Question Type", ["Multiple Choice", "Subjective", "Fill in the Blanks"])    
            difficulty_level = st.selectbox("Select Difficulty Level", ["Easy", "Medium", "Hard"])  
            num_questions = st.number_input("Number of Questions to Generate", min_value=1, value=10)  
            
            if st.button("Generate Questions"):    
                generated_questions = generate_questions_from_prompt(prompt, question_type, difficulty_level, num_questions, topic_name)    
                if generated_questions:    
                    st.write("Generated Questions:")    
                    selected_questions = []
                    for i, question in enumerate(generated_questions):    
                        if st.checkbox(f"Question {i+1}", value=True):
                            selected_questions.append(question)
                        st.write(question)    
                    
                    # Store generated questions in session state
                    if 'history' not in st.session_state:
                        st.session_state.history = []
                    st.session_state.history.append({"topic": topic_name, "questions": selected_questions})

                    st.session_state.generated_questions = selected_questions
                    st.session_state.topic_name = topic_name
                    st.success(f"Generated {len(selected_questions)} questions. Please proceed to add them to a question bank.")
                else:    
                    st.error("Failed to generate questions")        

            # New feature: View Generated Questions History
            st.subheader("View Generated Questions History")
            topic_questions = get_generated_questions_history()  # Retrieve the history

            if topic_questions:
                selected_topic = st.selectbox("Select Topic", options=list(topic_questions.keys()))
                if selected_topic:
                    questions = topic_questions[selected_topic]
                    st.write(f"Generated Questions for Topic: {selected_topic}")
                    for i, question in enumerate(questions, 1):
                        st.write(f"{i}. {question}")
            else:
                st.info("No generated questions history available.")

            if 'generated_questions' in st.session_state and st.session_state.generated_questions:
                st.subheader("Add Questions to Question Bank")
                
                existing_question_banks = get_all_question_banks()
                qb_options = ["Create New Question Bank"] + [f"ID: {str(qb['_id'])} - {qb['technology']} - {qb['difficulty']}" for qb in existing_question_banks]
                selected_qb = st.selectbox("Select Question Bank", options=qb_options)
                
                if st.button("Add Questions to Selected Bank"):
                    if selected_qb == "Create New Question Bank":
                        new_qb_technology = st.text_input("Enter technology for new question bank")
                        new_qb_difficulty = st.selectbox("Select difficulty for new question bank", ["Easy", "Medium", "Hard"])
                        if st.button("Create and Add"):
                            new_qb_id = create_new_question_bank(new_qb_technology, new_qb_difficulty, st.session_state.generated_questions)
                            if new_qb_id:
                                st.success(f"Created new question bank with ID: {new_qb_id} and added selected questions.")
                                st.session_state.generated_qb_id = new_qb_id  # Store the new question bank ID
                            else:
                                st.error("Failed to create new question bank.")
                    else:
                        qb_id = ObjectId(selected_qb.split('-')[0].split(':')[1].strip()) # Convert to ObjectId
                        if add_questions_to_question_bank(qb_id, st.session_state.topic_name, st.session_state.generated_questions):
                            st.success(f"Questions added to question bank ID: {qb_id}")
                        else:
                            st.error("Failed to add questions to the selected question bank.")

        # Clear the generated questions from session state
                        del st.session_state.generated_questions
                        del st.session_state.topic_name
        # Display content based on the selected option
    
    elif selected_tab == "Chatbot":
        # Display chatbot interface at the top
        st.subheader("Chat with the AI Trainer 🤖")
        
        # Initialize the message state if not exists
        if "msg" not in st.session_state:
            st.session_state.msg = ""
        
        # Create a container for the chat interface
        chat_container = st.container()

        # Define avatars
        user_avatar = "https://static.vecteezy.com/system/resources/previews/009/664/418/non_2x/people-user-team-transparent-free-png.png"
        ai_avatar = "https://thumbs.dreamstime.com/b/chatbot-logo-messenger-ai-robot-icon-vector-illustration-277900892.jpg"

        def clear_text():
            st.session_state.msg = st.session_state.user_input
            st.session_state.user_input = ""

        with chat_container:
            # Display chat messages
            for chat in st.session_state.chat_history:
                if chat['role'] == 'assistant':
                    # Chatbot message with avatar
                    st.markdown(
                        f"<div style='display: flex; align-items: center; margin: 5px 0;color:black'>"
                        f"<img src='{ai_avatar}' style='width: 40px; height: 40px; border-radius: 50%; margin-right: 10px;'>"
                        f"<div style='background-color: #e1ffc7; padding: 10px; border-radius: 10px; max-width: 80%;'>"
                        f"<strong>AI:</strong> {chat['content']}</div></div>",
                        unsafe_allow_html=True
                    )
                else:
                    # User message with avatar
                    st.markdown(
                        f"<div style='display: flex; align-items: center; margin: 5px 0; justify-content: flex-end;color:black'>"
                        f"<div style='background-color: #dcf8c6; padding: 10px; border-radius: 10px; max-width: 80%; margin-left: auto;'>"
                        f"<strong>You:</strong> {chat['content']}</div>"
                        f"<img src='{user_avatar}' style='width: 40px; height: 40px; border-radius: 50%; margin-left: 10px;'>"
                        f"</div>", 
                        unsafe_allow_html=True
                    )

            # Input field for user to enter a prompt
            st.text_input("Type your message here...", key="user_input", placeholder="Type a message...", on_change=clear_text)

            if st.session_state.msg:  # Only process if there's a message
                # Append user input to chat history
                st.session_state.chat_history.append({"role": "user", "content": st.session_state.msg})

                try:
                    # Generate AI response
                    prompt = f"You are an AI assistant for trainers. Respond to the following message: {st.session_state.msg}"
                    response = model.generate_content(prompt)
                    
                    # Handle the response properly for Gemini model
                    if hasattr(response, 'parts'):
                        ai_response = ''.join(part.text for part in response.parts)
                    else:
                        ai_response = response.candidates[0].content.parts[0].text
                    
                    # Append AI response to chat history
                    st.session_state.chat_history.append({"role": "assistant", "content": ai_response})
                except Exception as e:
                    st.error(f"Error generating response: {str(e)}")
                    ai_response = "I apologize, but I encountered an error. Please try again."
                    st.session_state.chat_history.append({"role": "assistant", "content": ai_response})
                
                # Clear the message state
                st.session_state.msg = ""
                
                # Rerun the app to display the new messages
                st.rerun()

                    
    
    

        
    notifications = get_notifications("trainer", None)  # Get notifications for trainer
    display_notifications(notifications, "trainer")  # Display notifications in the sidebar

    if notifications:  
        st.sidebar.write("Notifications:")  
        for notification in notifications:  
            st.sidebar.write(notification['message'])  
    else:  
        st.sidebar.write("No notifications available.")



    
# Download NLTK data
nltk.download('vader_lexicon')


def ascii_to_string(ascii_list):
    """Convert a list of ASCII values to a string."""
    return ''.join(chr(num) for num in ascii_list)

    # Example usage when retrieving topics
    curricula = get_all_curricula()
    for curriculum in curricula:
        # Assuming 'topics' is a list of ASCII values in the curricula
        if isinstance(curriculum['topics'], list):  # Check if topics is a list of ASCII values
            curriculum['topics'] = ascii_to_string(curriculum['topics'])  # Convert to string


def format_timestamp(timestamp_str):
    """
    Format timestamp for display
    """
    if timestamp_str == 'N/A':
        return 'N/A'
    try:
        # Parse the timestamp string to datetime
        timestamp = datetime.datetime.strptime(timestamp_str, "%Y-%m-%d %H:%M:%S")
        # Format it to a more readable format
        return timestamp.strftime("%b %d, %Y %I:%M %p")
    except (ValueError, TypeError):
        return 'N/A'



# Utility functions
def extract_text_from_file(file):
    file_extension = os.path.splitext(file.name)[1].lower()
    text = ""

    try:
        if file_extension == '.pdf':
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        elif file_extension == '.docx':
            doc = docx.Document(file)
            text = "\n".join([para.text for para in doc.paragraphs])
        elif file_extension == '.txt':
            text = file.getvalue().decode('utf-8')
        elif file_extension in ['.ppt', '.pptx']:
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text += shape.text + "\n"
        elif file_extension == '.csv':
            csv_data = pd.read_csv(file)
            text = csv_data.to_string(index=False)
        else:
            # For other file types, attempt to read as text
            try:
                text = file.getvalue().decode('utf-8')
            except UnicodeDecodeError:
                raise ValueError(f"Unable to extract text from {file_extension} file.")
    except Exception as e:
        raise ValueError(f"Error processing {file_extension} file: {str(e)}")

    # Clean the extracted text
    cleaned_text = clean_text(text)
    return cleaned_text

def clean_text(text):
    # Remove non-printable characters and control characters
    text = re.sub(r'[^\x20-\x7E]+', ' ', text)  # Keep only printable ASCII characters
    
    # Normalize whitespace
    text = re.sub(r'\s+', ' ', text)  # Replace multiple spaces with a single space
    text = text.strip()  # Trim leading and trailing whitespace

    return text


def create_new_question_bank(technology, difficulty, questions):
    db = create_connection()
    if db is None:
        return None

    try:
        questions_text = '\n'.join(questions)
        new_qb_doc = {
            "technology": technology,
            "difficulty": difficulty,
            "questions": questions_text,
            "created_at": datetime.now()
        }
        result = db.question_banks.insert_one(new_qb_doc)
        return str(result.inserted_id) # Return the ObjectId as a string
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None

def generate_questions(text, num_questions=5, question_type="multiple-choice"):
    if question_type == "multiple-choice":
        prompt = f"Generate {num_questions} multiple-choice questions based on the following text:\n\n{text}\n\nProvide the questions and options in the following format:\n\nQ1: [Question]\nA) [Option 1]\nB) [Option 2]\nC) [Option 3]\nD) [Option 4]\n\nQ2: [Question]\nA) [Option 1]\nB) [Option 2]\nC) [Option 3]\nD) [Option 4]\n\n..."
    elif question_type == "subjective":
        prompt = f"Generate {num_questions} subjective questions based on the following text:\n\n{text}\n\nProvide the questions in the following format:\n\nQ1: [Question]\n\nQ2: [Question]\n\n..."
    elif question_type == "fill-in-the-blank":
        prompt = f"Generate {num_questions} fill-in-the-blank questions based on the following text:\n\n{text}\n\nProvide the questions and correct answers in the following format:\n\nQ1: [Question]\nA: [Correct Answer]\n\nQ2: [Question]\nA: [Correct Answer]\n\n..."
    else:
        raise ValueError("Invalid question type")

    response = model.generate_content(prompt)
    generated_text = response.text

    questions = []
    options = []
    correct_answers = []

    lines = [line.strip() for line in generated_text.split('\n') if line.strip()]

    i = 0
    while i < len(lines):
        if lines[i].startswith('Q'):
            question = lines[i].split(': ', 1)[1]
            questions.append(question)
            if question_type == "multiple-choice":
                options_list = []
                correct_answer = None
                # Look for options immediately following the question
                for j in range(i + 1, len(lines)):
                    if lines[j].startswith(('A)', 'B)', 'C)', 'D)')):
                        option = lines[j].split(') ', 1)[1]
                        options_list.append(option)
                        if lines[j].startswith('A)'): # Assuming A is always the correct answer for simplicity in parsing
                            correct_answer = option
                    else:
                        break # Stop if a line doesn't start with an option letter
                options.append(options_list)
                correct_answers.append(correct_answer)
                i = j # Move index to the line after the last option processed
            elif question_type == "fill-in-the-blank":
                if i+1 < len(lines) and lines[i+1].startswith('A:'):
                    options.append([lines[i+1].split(': ', 1)[1]])
                    correct_answers.append(lines[i+1].split(': ', 1)[1])
                    i += 2
                else:
                    options.append([""])
                    correct_answers.append("")
                    i += 1
            else: # subjective
                options.append([])
                correct_answers.append("")
                i += 1
        else:
            i += 1

    return questions[:num_questions], options[:num_questions], correct_answers[:num_questions]

# Removed ensure_table_exists as MongoDB handles collection creation implicitly

def review_feedback():
    # Fetch feedback data from MongoDB
    db = create_connection()
    if db is None:
        return []
    
    try:
        feedback_data = list(db.feedback.find({}))
        # Convert ObjectId to string for compatibility with DataFrame
        for item in feedback_data:
            if '_id' in item:
                item['id'] = str(item['_id'])
                del item['_id'] # Remove ObjectId if not needed for display
            if 'question_bank_id' in item and isinstance(item['question_bank_id'], ObjectId):
                item['question_bank_id'] = str(item['question_bank_id'])
        return feedback_data
    except OperationFailure as e:
        st.error(f"Error retrieving feedback: {e}")
        return []

def analyze_sentiment(text):
    sia = SentimentIntensityAnalyzer()
    sentiment_score = sia.polarity_scores(text)['compound']
    if sentiment_score > 0.05:
        return 'Positive'
    elif sentiment_score < -0.05:
        return 'Negative'
    else:
        return 'Neutral'

def display_questions(questions, options, correct_answers):
    for i, question in enumerate(questions):
        st.write(question)
        if options[i]:
            st.write("Options:")
            for j, option in enumerate(options[i]):
                st.write(f"{chr(65+j)}) {option}")
        st.write(f"Correct Answer: {correct_answers[i]}")
        st.write("")
# User Authentication Functions
def login_user(username, password):
    db = create_connection()
    if db is None:
        return None

    user = db.users.find_one({"username": username})

    if user and check_password_hash(user['password'], password):
        # Convert ObjectId to string for session state
        user['_id'] = str(user['_id'])
        return user
    return None

def register_user(email, username, password, role):
    # Validate email format
    email_regex = r'^[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+$'
    if not re.match(email_regex, email):
        st.error("Invalid email format. Please enter a valid email address.")
        return False

    db = create_connection()
    if db is None:
        return False

    existing_user = db.users.find_one({"username": username})

    if existing_user:
        return False

    hashed_password = generate_password_hash(password)
    user_data = {
        "email": email,
        "username": username,
        "password": hashed_password,
        "role": role
    }
    try:
        db.users.insert_one(user_data)
        return True
    except OperationFailure as e:
        st.error(f"Database error during registration: {e}")
        return False
# Administrator Functions
def get_system_stats():
    db = create_connection()
    if db is None:
        return None

    stats = {}
    try:
        # Get count of documents in each collection
        stats['users'] = db.users.count_documents({})
        stats['question_banks'] = db.question_banks.count_documents({})
        stats['learning_plans'] = db.learning_plans.count_documents({})
        stats['feedback'] = db.feedback.count_documents({})
        return stats
    except OperationFailure as e:
        st.error(f"Error retrieving system stats: {e}")
        return None

def get_all_users():
    db = create_connection()
    if db is None:
        return []

    try:
        users_cursor = db.users.find({}, {"username": 1, "email": 1, "role": 1, "_id": 0})
        users = list(users_cursor)
        return users
    except OperationFailure as e:
        st.error(f"Error retrieving all users: {e}")
        return []

def update_user_role(username, new_role):
    db = create_connection()
    if db is None:
        return False

    try:
        user_exists = db.users.find_one({"username": username})
        if not user_exists:
            st.error("User does not exist.")
            return False

        result = db.users.update_one({"username": username}, {"$set": {"role": new_role}})
        return result.modified_count > 0
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False

# Trainer Functions
def upload_curriculum(technology, topics, content):   
  db = create_connection()   
  if db is None:   
    return False   
   
  try:   
    topics_str = ','.join(topics)   
    
    # Determine the content type (file-like object or string)   
    if hasattr(content, 'read'):   
      # File-like object   
      content_text = content.read().decode('utf-8')   
    else:   
      # String   
      content_text = content   
   
    # Generate questions from the content   
    questions, options, correct_answers = generate_questions(content_text)   
   
    # Convert questions, options and correct_answers to strings   
    questions_str = '|||'.join(questions)   
    options_str = '|||'.join(['###'.join(option) for option in options]) # Using '###' as separator for options   
    correct_answers_str = '|||'.join([','.join(map(str, ans)) if isinstance(ans, list) else str(ans) for ans in correct_answers])   
   
    # Insert or update the curriculum using upsert
    curriculum_doc = {
        "technology": technology,
        "topics": topics_str,
        "filename": "curriculum_" + technology + ".txt",
        "content": content_text
    }
    db.curriculum.update_one({"technology": technology}, {"$set": curriculum_doc}, upsert=True)
   
    # Insert the generated questions into the generated_question_files collection
    generated_questions_doc = {
        "technology": technology,
        "topics": topics_str,
        "questions": questions_str,
        "options": options_str,
        "correct_answers": correct_answers_str,
        "created_at": datetime.now()
    }
    db.generated_question_files.insert_one(generated_questions_doc)

    return True   
  except OperationFailure as err:   
    st.error(f"Database error: {err}")   
    return False   
  except Exception as e:
    st.error(f"Error in upload_curriculum: {e}")
    return False


def get_curriculum_text(technology):
    db = create_connection()
    if db is None:
        return None

    try:
        result = db.curriculum.find_one({"technology": technology}, {"topics": 1, "_id": 0})
        if result:
            return result.get('topics')
        else:
            st.error(f"No curriculum content found for technology: {technology}")
            return None
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None

def save_question_bank(technology, topics, questions, difficulty, correct_answers, question_type, options=None):  # ✅ Correct one

    db = create_connection()
    if db is None:
        return False

    try:
        # Prepare document for question_banks collection
        qb_doc = {
            "technology": technology,
            "topics": topics, # topics is already a string
            "questions": questions, # questions is already a single string
            "difficulty": difficulty,
            "question_type": question_type,
            "options": options, # options is already a single string
            "created_at": datetime.now()
        }
        
        # Insert into question_banks collection
        result_qb = db.question_banks.insert_one(qb_doc)
        question_bank_id = result_qb.inserted_id # MongoDB's _id

        # Prepare document for question_answers collection
        answer_doc = {
            "question_bank_id": question_bank_id, # Link to the question bank
            "answer_data": correct_answers # correct_answers is already a single string
        }
        
        # Insert into question_answers collection
        db.question_answers.insert_one(answer_doc)
        
        return str(question_bank_id) # Return as string for consistency with app logic
        
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False
    except Exception as e:
        st.error(f"General error in save_question_bank: {e}")
        return False

def get_topics_for_technology(technology):
    db = create_connection()
    if db is None:
        return None

    try:
        result = db.curriculum.find_one({"technology": technology}, {"topics": 1, "_id": 0})
        if result and 'topics' in result:
            return result['topics'].split(',')
        else:
            st.error(f"No topics found for technology: {technology}")
            return None
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None

def get_all_curricula():
    db = create_connection()
    if db is None:
        return None

    try:
        curricula_cursor = db.curriculum.find({}, {"technology": 1, "topics": 1}) # Include _id for potential future use
        curricula = []
        for doc in curricula_cursor:
            doc['id'] = str(doc['_id']) # Add string version of _id
            curricula.append(doc)
        return curricula
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None

def get_all_question_banks():
    db = create_connection()
    if db is None:
        return None

    try:
        # Fetch all documents from the question_banks collection
        question_banks_cursor = db.question_banks.find({})
        question_banks = list(question_banks_cursor)
        return question_banks
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None

# Employee Functions
def get_learning_plan(username):
    db = create_connection()
    if db is None:
        return None

    learning_plan = db.learning_plans.find_one({"username": username})
    return learning_plan

def submit_feedback(username, question_bank_id, feedback_text, rating, feedback_type):
    sentiment = analyze_sentiment(feedback_text)  # Optional: Analyze sentiment
    db = create_connection()
    if db is None:
        return False

    try:
        feedback_doc = {
            "username": username,
            "question_bank_id": ObjectId(question_bank_id) if question_bank_id else None, # Store as ObjectId
            "feedback_text": feedback_text,
            "rating": rating,
            "sentiment": sentiment,
            "feedback_type": feedback_type,
            "created_at": datetime.now()
        }
        db.feedback.insert_one(feedback_doc)

        # Prepare notification message
        feedback_summary = f"New feedback received from {username}. Type: {feedback_type}. Rating: {rating}. Feedback: {feedback_text}"

        # Send notifications based on feedback type
        if feedback_type in ["User  Experience"] and rating >= 3:
            send_notification("admin", feedback_summary, username)
        elif feedback_type in ["Question Bank Feedback", "Assessment Feedback"]:
            send_notification("admin", feedback_summary, username)
            send_notification("trainer", feedback_summary, username)

        return True
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False


def take_assessment():
    st.subheader("Take Assessment")
    question_banks = get_all_question_banks()
    if not question_banks:
        st.info("No question banks available yet.")
    else:
        selected_qb = st.selectbox(
            "Select Question Bank",
            options=[(str(qb['_id']), f"{qb['technology']} - {qb['difficulty']}") for qb in question_banks],
            format_func=lambda x: f"ID: {x[0]} - {x[1]}",
            key="take_assessment_qb_select"
        )

        if selected_qb:
            qb_id_str, _ = selected_qb
            qb_id = ObjectId(qb_id_str) # Convert to ObjectId
            qb_details = next((qb for qb in question_banks if qb['_id'] == qb_id), None)
            if qb_details:
                questions = qb_details.get('questions', '').split('\n')
                options = qb_details.get('options', '').split('###') if qb_details.get('options') else []
                correct_answers_str = get_correct_answers(qb_id) # This returns a list of strings
                
                question_type = qb_details.get('question_type')

                score = 0
                user_answers = []
                for i, question in enumerate(questions):
                    if not question.strip(): # Skip empty questions
                        continue

                    st.write(f"**Q{i+1}:** {question.strip()}")
                    
                    if question_type == "multiple-choice" and options:
                        # Assuming options are stored as 'Opt1###Opt2###Opt3###Opt4' per question
                        # Need to parse options for the current question correctly
                        # For now, let's assume options are structured to match questions
                        if i < len(options): # Ensure we have options for this question
                            current_options = options[i].split('###')
                            if current_options:
                                answer = st.radio("Select an option", current_options, key=f"question_{i}")
                                user_answers.append(answer)
                            else:
                                user_answers.append("") # No options, no answer
                        else:
                            user_answers.append("") # No options for this question
                    elif question_type == "fill-in-the-blank":
                        answer = st.text_input("Enter your answer", key=f"question_{i}")
                        user_answers.append(answer)
                    elif question_type == "subjective":
                        answer = st.text_area("Enter your answer", key=f"question_{i}")
                        user_answers.append(answer)
                    else:
                        user_answers.append("") # For cases where question_type is not set or options are missing

                if st.button("Submit"):
                    total_questions_answered = 0
                    correct_answers_count = 0
                    for i, user_answer in enumerate(user_answers):
                        if i < len(correct_answers_str) and user_answer.strip().lower() == correct_answers_str[i].strip().lower():
                            st.success(f"Q{i+1}: Correct!")
                            correct_answers_count += 1
                        elif i < len(correct_answers_str):
                            st.error(f"Q{i+1}: Incorrect. Correct answer: {correct_answers_str[i].strip()}")
                        total_questions_answered += 1

                    st.write(f"Your score is {correct_answers_count} out of {total_questions_answered}")

                    # Save the assessment result
                    save_assessment_result(st.session_state.user['username'], qb_id, correct_answers_count)

def get_available_question_banks(username):
    db = create_connection()
    if db is None:
        return []

    try:
        # Get technologies from learning plans for the user
        learning_plans_cursor = db.learning_plans.find({"username": username}, {"technology": 1, "_id": 0})
        technologies = [lp['technology'] for lp in learning_plans_cursor if 'technology' in lp]

        # Get IDs of question banks already completed by the user
        completed_assessments_cursor = db.assessments.find({"username": username}, {"question_bank_id": 1, "_id": 0})
        completed_qb_ids = [a['question_bank_id'] for a in completed_assessments_cursor if 'question_bank_id' in a]

        # Find question banks matching learning plan technologies and not yet completed
        query = {
            "technology": {"$in": technologies},
            "_id": {"$nin": completed_qb_ids}
        }
        question_banks_cursor = db.question_banks.find(query, {"technology": 1, "topics": 1})
        
        # Convert ObjectId to string for 'id' field in the returned dictionary
        question_banks = []
        for qb in question_banks_cursor:
            qb['id'] = str(qb['_id'])
            question_banks.append(qb)
        
        return question_banks
    except OperationFailure as e:
        st.error(f"Error retrieving available question banks: {e}")
        return []

def get_completed_assessments(username):
    db = create_connection()
    if db is None:
        return []

    try:
        # Fetch assessments for the user
        assessments_cursor = db.assessments.find({"username": username}).sort("completed_at", -1)
        completed_assessments = []

        for assessment in assessments_cursor:
            # Fetch corresponding question bank details
            qb_details = db.question_banks.find_one({"_id": assessment['question_bank_id']})
            
            if qb_details:
                total_questions = len(qb_details.get('questions', '').split('\n')) if qb_details.get('questions') else 0
                
                percentage = round((assessment['score'] / total_questions) * 100, 2) if total_questions > 0 else 0.0

                completed_assessments.append({
                    'id': str(assessment['_id']), # Convert ObjectId to string
                    'question_bank_id': str(assessment['question_bank_id']), # Convert ObjectId to string
                    'technology': qb_details.get('technology'),
                    'difficulty': qb_details.get('difficulty'),
                    'score': assessment.get('score'),
                    'question_type': qb_details.get('question_type'),
                    'completed_at': assessment.get('completed_at').strftime('%Y-%m-%d %H:%M:%S') if assessment.get('completed_at') else 'N/A',
                    'total_questions': total_questions,
                    'percentage': percentage
                })
        return completed_assessments
    except OperationFailure as e:
        st.error(f"Error retrieving completed assessments: {e}")
        return []


def admin_dashboard():  
      
    with st.sidebar:
      # Create a sidebar for navigation using option_menu
      selected_tab = option_menu(
          menu_title="Admin Dashboard",  # required
          options=["System Stats", 
                  "User  Management", 
                  "Reports", 
                  "Employee Performance"],  # required
          icons=["bar-chart", "people", "file-earmark-text", "person-check"],  # optional
          menu_icon="cast",  # optional
          default_index=0,  # optional
          orientation="vertical",
      )

    # Display the selected tab content
    if selected_tab == "System Stats":
        st.subheader("System Statistics 📊")  
        stats = get_system_stats()  
        if stats:  
            st.write(f"Total Users: {stats.get('users', 0)}")  
            st.write(f"Total Question Banks: {stats.get('question_banks', 0)}")  
            st.write(f"Total Learning Plans: {stats.get('learning_plans', 0)}")  
            st.write(f"Total Feedback Entries: {stats.get('feedback', 0)}")  

            # Add system details  
            st.subheader("System Details 🖥️")  
            st.write(f"Operating System: {platform.system()}")  
            st.write(f"Platform: {platform.platform()}")  
            st.write(f"Processor: {platform.processor()}")  
            st.write(f"Python Version: {platform.python_version()}")  
        else:  
            st.error("Failed to retrieve system statistics") 

    elif selected_tab == "User  Management":
        st.subheader("User  Management �")
        
        # Check if users are already loaded in session state
        if 'users' not in st.session_state:
            st.session_state.users = get_all_users()  # Load users for the first time

        # Create a layout for the refresh button
        col1, col2 = st.columns([4, 1])  # Adjust column widths

        with col1:
            st.write("")  # Empty space for alignment

        with col2:
            # Refresh button with an icon
            if st.button("🔄", key="refresh_users", help="Refresh User List"):
                st.session_state.users = get_all_users()  # Refresh the user list

        users = st.session_state.users  # Use the loaded users from session state
        
        if users:
            # Display the user table
            user_table = []
            for user in users:
                user_table.append({
                    'Username': user.get('username'),
                    'Email': user.get('email'),
                    'Role': user.get('role')
                })
            
            user_df = pd.DataFrame(user_table)
            st.table(user_df)

            # CSS to inject for compact layout
            st.markdown("""
                <style>
                .stSelectbox {
                    margin-bottom: 0px;
                }
                .stButton {
                    display: inline-block;
                    margin-right: 10px;
                }
                .user-row {
                    margin-bottom: 10px;
                }
                </style>
            """, unsafe_allow_html=True)

            # Search feature
            st.subheader("Search User")
            search_username = st.text_input("Enter username to search:", "")
            
            if search_username:
                filtered_users = [user for user in users if search_username.lower() in user.get('username', '').lower()]
            else:
                filtered_users = users

            # Display users (filtered or all)
            if filtered_users:
                for user in filtered_users:
                    with st.container():
                        cols = st.columns([2, 3, 3])  # Adjusted column widths
                        
                        with cols[0]:  # Username column
                            st.write(user.get('username'))
                        
                        with cols[1]:  # Role selection column
                            new_role = st.selectbox(
                                f"New Role for {user.get('username')}", 
                                ["None", "Administrator", "Trainer", "Employee"], 
                                key=f"new_role_{user.get('username')}",
                                label_visibility="collapsed"  # Hides the label
                            )
                        
                        with cols[2]:  # Buttons column
                            c1, c2 = st.columns([1, 1])  # Split the column for buttons
                            with c1:
                                if st.button("Update Role", key=f"update_role_button_{user.get('username')}", use_container_width=True):
                                    if update_user_role(user.get('username'), new_role):
                                        st.success(f"Role updated for {user.get('username')}")
                                        # Refresh the user list after update
                                        st.session_state.users = get_all_users()
                                    else:
                                        st.error("Failed to update role")
                            
                            with c2:
                                if st.button("Remove User", key=f"remove_user_button_{user.get('username')}", use_container_width=True):
                                    if remove_user(user.get('username')):
                                        st.success(f"User {user.get('username')} removed successfully")
                                        # Refresh the user list after removal
                                        st.session_state.users = get_all_users()
                                    else:
                                        st.error("Failed to remove user")
            else:
                st.info("No users found with that username.")
        else:
            st.info("No users available.")

    elif selected_tab == "Reports":
        st.subheader("Generate Reports 🔍")
        report_type = st.selectbox("Select Report Type", 
                                    ["User Activity", "Question Bank Usage", 
                                    "Feedback Summary", "Sentiment Analysis", 
                                    "Employee Performance"])

        if st.button("Generate Report"):
            if report_type == "User Activity":
                user_activity_report()
            elif report_type == "Question Bank Usage":
                question_bank_usage_report()
            elif report_type == "Feedback Summary":
                feedback_summary_report()
            elif report_type == "Sentiment Analysis":
                sentiment_analysis_report()
            elif report_type == "Employee Performance":
                employee_performance_report()



    elif selected_tab == "Employee Performance":
        st.subheader("Employee Performance 🎯")
        employees = get_all_users()

        if employees:
            selected_employee = st.selectbox(
                "Select Employee",
                options=[employee['username'] for employee in employees],
                key="employee_performance_select"
            )

            if selected_employee:
                # Fetch assessment results for the selected employee
                assessment_results = get_assessment_results(selected_employee)
                if assessment_results:
                    # Prepare data for the table
                    performance_data = []
                    for result in assessment_results:
                        performance_data.append({
                            'Question Bank ID': str(result['question_bank_id']), # Convert ObjectId to string
                            'Score': result['score'],
                            'Completed At': result['completed_at']
                        })

                    # Convert to DataFrame for better visualization
                    performance_df = pd.DataFrame(performance_data)

                    # Display summary metrics
                    st.subheader(f"Summary Statistics for {selected_employee}")
                    total_assessments = len(performance_df)
                    avg_score = performance_df['Score'].mean() if total_assessments > 0 else 0
                    best_score = performance_df['Score'].max() if total_assessments > 0 else 0

                    col1, col2, col3 = st.columns(3)
                    with col1:
                        st.metric("Total Assessments", total_assessments)
                    with col2:
                        st.metric("Average Score", f"{avg_score:.1f}")
                    with col3:
                        st.metric("Best Score", best_score)

                    # Display the performance data in a styled table
                    st.write(f"Performance Data for {selected_employee}:")
                    st.dataframe(performance_df.style.highlight_max(axis=0))  # Highlight max scores

                    # Convert 'Completed At' column to datetime for sorting
                    performance_df['Completed At'] = pd.to_datetime(performance_df['Completed At'])
                    performance_df.sort_values('Completed At', inplace=True)

                    # Create visualizations
                    st.subheader("Performance Over Time")

                    # Line chart for scores over time
                    fig_line = px.line(performance_df, x='Completed At', y='Score', 
                                    title='Score Over Time', markers=True)
                    st.plotly_chart(fig_line)

                    # Bar chart for scores by question bank
                    fig_bar = px.bar(performance_df, x='Question Bank ID', y='Score', 
                                    title='Scores by Question Bank', text='Score')
                    st.plotly_chart(fig_bar)

                    # Convert figures to HTML format for download
                    fig_line_html = fig_line.to_html(full_html=False)
                    fig_bar_html = fig_bar.to_html(full_html=False)

                    # Provide download buttons
                    st.download_button(label="Download Line Chart as HTML", data=fig_line_html, 
                                    file_name=f"{selected_employee}_performance_over_time.html", mime="text/html")

                    st.download_button(label="Download Bar Chart as HTML", data=fig_bar_html, 
                                    file_name=f"{selected_employee}_score_by_question_bank.html", mime="text/html")

                    st.download_button(label="Download Performance Data as CSV", 
                                    data=performance_df.to_csv(index=False), 
                                    file_name=f"{selected_employee}_performance.csv", mime="text/csv")

                else:
                    st.info("No assessment results available for this employee.")
        else:
            st.info("No employees available.")

    # Display notifications in the sidebar
     
    notifications = get_notifications("admin", None)  # Get notifications for admin
    display_notifications(notifications, "admin")  # Display notifications in the sidebar
    if notifications:  
        st.sidebar.write("Notifications:")  
        for notification in notifications:  
            st.sidebar.write(notification['message'])  
    else:  
        st.sidebar.write("No notifications available.")
    

import streamlit as st
import plotly.express as px
import pandas as pd
import io

def employee_performance_report():
    db = create_connection()
    if db is None:
        st.error("Failed to connect to database")
        return

    try:
        # Aggregation pipeline to get employee performance
        pipeline = [
            {"$group": {
                "_id": "$username",
                "num_assessments": {"$sum": 1},
                "avg_score": {"$avg": "$score"}
            }},
            {"$project": {
                "username": "$_id",
                "num_assessments": 1,
                "avg_score": 1,
                "_id": 0
            }}
        ]
        results = list(db.assessments.aggregate(pipeline))

        if results:
            df = pd.DataFrame(results)

            # Create an interactive bar chart using Plotly
            fig = px.bar(df, x='username', y='avg_score', title='Employee Performance',
                         labels={'username': 'Employee', 'avg_score': 'Average Score'},
                         color='avg_score', color_continuous_scale=px.colors.sequential.Viridis)

            # Display the chart in Streamlit
            st.plotly_chart(fig)

            # Convert figure to HTML format for download
            fig_html = fig.to_html(full_html=False)

            # Provide a download button for the interactive HTML chart
            st.download_button(
                label="Download Chart as HTML",
                data=fig_html,
                file_name="employee_performance_chart.html",
                mime="text/html"
            )

            # Download CSV Data
            csv = df.to_csv(index=False)
            st.download_button(label="Download as CSV", data=csv, file_name="employee_performance.csv", mime="text/csv")

        else:
            st.error("No employee performance data available")
    except OperationFailure as e:
        st.error(f"Database error during report generation: {e}")


      

def user_activity_report():
    db = create_connection()
    if db is None:
        st.error("Failed to connect to database")
        return

    try:
        # Aggregation pipeline to get user activity for employees
        pipeline = [
            {"$lookup": {
                "from": "users",
                "localField": "username",
                "foreignField": "username",
                "as": "user_info"
            }},
            {"$unwind": "$user_info"},
            {"$match": {"user_info.role": "Employee"}},
            {"$group": {
                "_id": "$username",
                "num_assessments": {"$sum": 1},
                "avg_score": {"$avg": "$score"}
            }},
            {"$project": {
                "username": "$_id",
                "num_assessments": 1,
                "avg_score": 1,
                "_id": 0
            }}
        ]
        results = list(db.assessments.aggregate(pipeline))

        if results:
            df = pd.DataFrame(results)

            # Visualization - Pie Chart
            fig = px.pie(df, values='num_assessments', names='username', title='User Activity Distribution')
            st.plotly_chart(fig)

            # Convert figure to HTML format for download
            fig_html = fig.to_html(full_html=False)

            # Download buttons
            st.download_button(label="Download Chart as HTML", data=fig_html, file_name="user_activity_chart.html", mime="text/html")
            st.download_button(label="Download as CSV", data=df.to_csv(index=False), file_name="user_activity.csv", mime="text/csv")

        else:
            st.error("No user activity data available")
    except OperationFailure as e:
        st.error(f"Database error during report generation: {e}")


def remove_user(username):   
    db = create_connection()   
    if db is None:   
        return False   

    try:   
        # First, delete related records in the assessments collection
        db.assessments.delete_many({"username": username})   

        # Now, delete the user from the users collection
        result = db.users.delete_one({"username": username})   
        
        return result.deleted_count > 0
    except OperationFailure as err:   
        st.error(f"Database error: {err}")   
        return False   
  
def question_bank_usage_report():
    db = create_connection()
    if db is None:
        st.error("Failed to connect to database")
        return

    try:
        # Aggregation pipeline to get question bank usage
        pipeline = [
            {"$lookup": {
                "from": "question_banks",
                "localField": "question_bank_id",
                "foreignField": "_id",
                "as": "qb_info"
            }},
            {"$unwind": "$qb_info"},
            {"$group": {
                "_id": "$qb_info.technology",
                "num_assessments": {"$sum": 1}
            }},
            {"$project": {
                "technology": "$_id",
                "num_assessments": 1,
                "_id": 0
            }}
        ]
        results = list(db.assessments.aggregate(pipeline))

        if results:
            df = pd.DataFrame(results)

            # Visualization - Bar Chart
            fig = px.bar(df, x='technology', y='num_assessments', title='Question Bank Usage',
                         labels={'technology': 'Technology', 'num_assessments': 'Number of Assessments'},
                         color='num_assessments', color_continuous_scale=px.colors.sequential.Plasma)
            
            st.plotly_chart(fig)

            # Convert figure to HTML format for download
            fig_html = fig.to_html(full_html=False)

            # Download buttons
            st.download_button(label="Download Chart as HTML", data=fig_html, file_name="question_bank_usage_chart.html", mime="text/html")
            st.download_button(label="Download as CSV", data=df.to_csv(index=False), file_name="question_bank_usage.csv", mime="text/csv")

        else:
            st.error("No question bank usage data available")
    except OperationFailure as e:
        st.error(f"Database error during report generation: {e}")

  
def sentiment_analysis_report():
    db = create_connection()
    if db is None:
        st.error("Failed to connect to database")
        return

    try:
        # Aggregation pipeline to get sentiment counts
        pipeline = [
            {"$group": {
                "_id": "$sentiment",
                "num_feedback": {"$sum": 1}
            }},
            {"$project": {
                "sentiment": "$_id",
                "num_feedback": 1,
                "_id": 0
            }}
        ]
        results = list(db.feedback.aggregate(pipeline))

        if results:
            df = pd.DataFrame(results)

            # Visualization - Pie Chart
            fig = px.pie(df, values='num_feedback', names='sentiment', title='Sentiment Analysis',
                         color='sentiment', color_discrete_sequence=px.colors.qualitative.Set1)
            
            st.plotly_chart(fig)

            # Convert figure to HTML format for download
            fig_html = fig.to_html(full_html=False)

            # Download buttons
            st.download_button(label="Download Chart as HTML", data=fig_html, file_name="sentiment_analysis_chart.html", mime="text/html")
            st.download_button(label="Download as CSV", data=df.to_csv(index=False), file_name="sentiment_analysis.csv", mime="text/csv")

        else:
            st.error("No sentiment analysis data available")
    except OperationFailure as e:
        st.error(f"Database error during report generation: {e}")


def notify_trainer_of_bad_feedback(feedback_summary):
    send_notification("trainer", feedback_summary, "Admin Notification") 
  
def feedback_summary_report():
    db = create_connection()
    if db is None:
        st.error("Failed to connect to database")
        return

    try:
        # Aggregation pipeline to get feedback summary
        pipeline = [
            {"$group": {
                "_id": "$feedback_type",
                "avg_rating": {"$avg": "$rating"},
                "num_feedback": {"$sum": 1}
            }},
            {"$project": {
                "feedback_type": "$_id",
                "avg_rating": 1,
                "num_feedback": 1,
                "_id": 0
            }}
        ]
        results = list(db.feedback.aggregate(pipeline))

        if results:
            df = pd.DataFrame(results)

            # Visualization - Bar Chart
            fig = px.bar(df, x='feedback_type', y='num_feedback', title='Feedback Summary',
                         labels={'feedback_type': 'Feedback Type', 'num_feedback': 'Number of Feedback'},
                         color='num_feedback', color_continuous_scale=px.colors.sequential.Viridis)

            st.plotly_chart(fig)

            # Convert figure to HTML format for download
            fig_html = fig.to_html(full_html=False)

            # Download buttons
            st.download_button(label="Download Chart as HTML", data=fig_html, file_name="feedback_summary_chart.html", mime="text/html")
            st.download_button(label="Download as CSV", data=df.to_csv(index=False), file_name="feedback_summary.csv", mime="text/csv")

        else:
            st.error("No feedback data available")
    except OperationFailure as e:
        st.error(f"Database error during report generation: {e}")

def get_bad_feedback():
    db = create_connection()
    if db is None:
        return []

    try:
        feedback_records = list(db.feedback.find({"rating": {"$lt": 3}}, {"username": 1, "feedback_text": 1, "rating": 1, "_id": 0}))
        return feedback_records
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return []

def feedback_form(username):
    st.subheader("Submit Feedback")
    
    # Select the type of feedback
    feedback_type = st.selectbox("Select Feedback Type:", 
                                options=["User  Experience", "Question Bank Feedback", "Assessment Feedback"])
    
    # Select a question bank for feedback if applicable
    question_banks = get_all_question_banks()
    selected_qb = None
    if feedback_type in ["Question Bank Feedback", "Assessment Feedback"] and question_banks:
        selected_qb = st.selectbox(
            "Select Question Bank for Feedback",
            options=[(str(qb['_id']), f"{qb['technology']} - {qb['difficulty']}") for qb in question_banks],
            format_func=lambda x: f"ID: {x[0]} - {x[1]}",
            key="feedback_qb_select"
        )

    feedback_text = st.text_area("Submit your feedback:", height=100)
    rating = st.slider("Rate your experience (1-5):", 1, 5, 3)

    if st.button("Submit Feedback"):
        question_bank_id = selected_qb[0] if selected_qb else None
        if submit_feedback(username, question_bank_id, feedback_text, rating, feedback_type):
            st.success("Feedback submitted successfully!")
        else:
            st.error("Failed to submit feedback.")

def feedback_received(feedback):  
   db = create_connection()  
   if db is None:  
      return False  
  
   try:  
      feedback_doc = {"feedback": feedback, "created_at": datetime.now()}
      db.feedback.insert_one(feedback_doc)
  
      print("Feedback received!")  
  
      # Send notification to admin  
      recipient_role = "admin"  
      message = f"New feedback received: {feedback}"  
      print(f"Sending notification to {recipient_role}: {message}")  
      send_notification(recipient_role, message, "system")  # Changed username to "system"
  
      # Send notification to trainer  
      recipient_role = "trainer"  
      message = f"New feedback received: {feedback}"  
      print(f"Sending notification to {recipient_role}: {message}")  
      send_notification(recipient_role, message, "system")  # Changed username to "system"
  
      # Send notification to employee  
      recipient_role = "employee"  
      message = f"New feedback received: {feedback}"  
      print(f"Sending notification to {recipient_role}: {message}")  
      send_notification(recipient_role, message, "system")  # Changed username to "system"
  
      return True  
   except OperationFailure as err:  
      st.error(f"Database error: {err}")  
      return False


    

def generate_pdf_document(questions):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)

    # Add questions
    for i, question in enumerate(questions, 1):
        pdf.multi_cell(0, 10, txt=f"{i}. {question}", align='L')
        pdf.ln(5)

    # Save to a temporary file
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_file:
        pdf.output(temp_file.name)  # Save the PDF to a temporary file
        temp_file.seek(0)  # Move to the beginning of the file
        buffer = temp_file.read()  # Read the content of the file into a buffer

    return buffer

def generate_docx_document(questions):
    doc = docx.Document()
    
    # Add title
    doc.add_heading('Question Bank', 0)
    
    # Add questions
    for i, question in enumerate(questions, 1):
        doc.add_paragraph(f"{i}. {question}")
        doc.add_paragraph()  # Add spacing
    
    # Save to buffer
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()

def generate_pptx_document(questions):
    prs = Presentation()
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = "Question Bank"
    
    # Add questions (one per slide)
    for i, question in enumerate(questions, 1):
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        title_shape.text = f'Question {i}'
        
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = question
    
    # Save to buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


def generate_csv_document(questions):
    # Create CSV in memory
  buffer = io.StringIO()
  for i, question in enumerate(questions, 1):
        buffer.write(f"{i},{question}\n")
    
  return buffer.getvalue().encode()

def add_questions_to_question_bank(qb_id, topic_name, new_questions):
    db = create_connection()
    if db is None:
        return False

    try:
        # Retrieve existing questions from the question bank
        result = db.question_banks.find_one({"_id": qb_id}, {"questions": 1})

        if result and 'questions' in result:
            existing_questions = result['questions'].split('\n')  # Split existing questions into a list
            updated_questions = existing_questions + new_questions  # Append new questions
            updated_questions_str = '\n'.join(updated_questions)  # Join the updated list into a string

            # Update the question bank with the new list of questions
            update_result = db.question_banks.update_one(
                {"_id": qb_id},
                {"$set": {"questions": updated_questions_str}}
            )
            return update_result.modified_count > 0
        else:
            return False
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False

def save_generated_questions(questions, topic_name):   
    db = create_connection()   
    if db is None:   
        return False   
   
    try:   
        generated_doc = {
            "topic_name": topic_name,
            "questions": json.dumps(questions),
            "created_at": datetime.now()
        }
        db.generated_questions_history.insert_one(generated_doc)
        return True   
    except OperationFailure as err:   
        st.error(f"Database error: {err}")   
        return False   
  
def get_generated_questions_history():
    db = create_connection()
    if db is None:
        return None

    try:
        results_cursor = db.generated_questions_history.find({}, {"topic_name": 1, "questions": 1, "_id": 0})
        topic_questions = {}
        for result in results_cursor:
            topic_name = result.get('topic_name')
            questions = json.loads(result.get('questions', '[]'))
            if topic_name:
                topic_questions[topic_name] = questions
        return topic_questions
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None
def generate_questions_from_prompt(prompt, question_type, difficulty_level, num_questions, topic_name):    
    try:    
        response = model.generate_content(f"Generate {num_questions} {question_type} questions based on the following text with {difficulty_level} difficulty:\n\n{prompt}")    
        generated_text = response.text    
        questions = [line.strip() for line in generated_text.split('\n') if line.strip()]    
        
        # Save generated questions to history
        save_generated_questions(questions, topic_name)  # Call to save the generated questions
        
        return questions    
    except Exception as e:    
        st.error(f"Error generating questions: {e}")    
        return None
def get_curriculum_history():  
  db = create_connection()  
  if db is None:  
    return None  
  
  try:
    results_cursor = db.curriculum.find({}, {"technology": 1, "topics": 1})
    results = []
    for doc in results_cursor:
        doc['id'] = str(doc['_id']) # Convert ObjectId to string for 'id'
        results.append(doc)
    return results
  except OperationFailure as err:
    st.error(f"Database error: {err}")
    return None

def generate_question_bank_document(qb_id):  
   db = create_connection()  
   if db is None:  
      return None  
  
   try:
      result = db.question_banks.find_one({"_id": ObjectId(qb_id)}, {"questions": 1})
  
      if result and 'questions' in result:  
         questions = result['questions'].split('\n')  
         doc = docx.Document()  
         for question in questions:  
           doc.add_paragraph(question)  
         buffer = io.BytesIO()  
         doc.save(buffer)  
         buffer.seek(0)  
         return buffer.getvalue()  
      else:  
         return None  
   except OperationFailure as err:
      st.error(f"Database error: {err}")
      return None
  
def generate_personalized_learning_plan_document(username):   
  db = create_connection()   
  if db is None:   
    return None   
   
  try:
    # Find learning plans for the user
    learning_plan_doc = db.learning_plans.find_one({"username": username})
    
    if learning_plan_doc:
        qb_id = learning_plan_doc.get('question_bank_id') # Assuming learning_plan stores qb_id
        if qb_id:
            learning_plan = prepare_learning_plan(qb_id, username)   
            if learning_plan:   
                doc = docx.Document()   
                doc.add_paragraph(f"Technology: {learning_plan['technology']}")   
                doc.add_paragraph(f"Start Date: {learning_plan['start_date']}")   
                doc.add_paragraph(f"End Date: {learning_plan['end_date']}")   
                doc.add_paragraph(f"Status: {learning_plan['status']}")   
                doc.add_paragraph(f"Estimated Time: {learning_plan['estimated_time']} hours")   
            
                buffer = io.BytesIO()   
                doc.save(buffer)   
                buffer.seek(0)   
                return buffer.getvalue()   
            else:   
                return None   
        else:
            return None # No question_bank_id in learning plan
    else:   
        return None
  except OperationFailure as err:
    st.error(f"Database error: {err}")
    return None
  
def get_all_feedback():
    db = create_connection()
    if db is None:
        return []

    try:
        feedback_records_cursor = db.feedback.find({}, {"username": 1, "question_bank_id": 1, "feedback_text": 1, "rating": 1, "sentiment": 1, "feedback_type": 1, "_id": 0})
        feedback_records = list(feedback_records_cursor)
        # Convert ObjectId to string for question_bank_id if it exists
        for record in feedback_records:
            if 'question_bank_id' in record and isinstance(record['question_bank_id'], ObjectId):
                record['question_bank_id'] = str(record['question_bank_id'])
        return feedback_records
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return []

def generate_feedback_report_document(qb_id):  
   db = create_connection()  
   if db is None:  
      return None  
  
   try:
      results_cursor = db.feedback.find({"question_bank_id": ObjectId(qb_id)}, {"feedback_text": 1, "rating": 1, "_id": 0})
      results = list(results_cursor)
  
      if results:  
         doc = docx.Document()  
         for result in results:  
           doc.add_paragraph(result.get('feedback_text', ''))  
           doc.add_paragraph(str(result.get('rating', 'N/A')))  
         buffer = io.BytesIO()  
         doc.save(buffer)  
         buffer.seek(0)  
         return buffer.getvalue()  
      else:  
         return None  
   except OperationFailure as err:
      st.error(f"Database error: {err}")
      return None
  
def generate_curriculum_mapping_document(technology):
    db = create_connection()
    if db is None:
        return None

    try:
        result = db.curriculum.find_one({"technology": technology}, {"topics": 1, "_id": 0})

        if result and 'topics' in result:
            topics = result['topics']
            # topics is already a string in MongoDB, no need for bytes decoding
            topics = topics.split(',')  # Now split the string into a list
            doc = docx.Document()
            for topic in topics:
                doc.add_paragraph(topic)
            buffer = io.BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            return buffer.getvalue()
        else:
            return None
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None
  
def generate_issue_resolution_report_document(qb_id):  
   db = create_connection()  
   if db is None:  
      return None  
  
   try:
      results_cursor = db.issue_resolution.find({"question_bank_id": ObjectId(qb_id)}, {"issue": 1, "resolution": 1, "_id": 0})
      results = list(results_cursor)
  
      if results:  
         doc = docx.Document()  
         for result in results:  
           doc.add_paragraph(result.get('issue', ''))  
           doc.add_paragraph(result.get('resolution', ''))  
         buffer = io.BytesIO()  
         doc.save(buffer)  
         buffer.seek(0)  
         return buffer.getvalue()  
      else:  
         return None  
   except OperationFailure as err:
      st.error(f"Database error: {err}")
      return None
  
def generate_assessment_completion_report_document(username):  
   db = create_connection()  
   if db is None:  
      return None  
  
   try:
      results_cursor = db.assessments.find({"username": username}, {"score": 1, "completed_at": 1, "_id": 0})
      results = list(results_cursor)
  
      if results:  
         doc = docx.Document()  
         for result in results:  
           doc.add_paragraph(str(result.get('score', 'N/A')))  
           doc.add_paragraph(str(result.get('completed_at', 'N/A')))  
         buffer = io.BytesIO()  
         doc.save(buffer)  
         buffer.seek(0)  
         return buffer.getvalue()  
      else:  
         return None
   except OperationFailure as err:
      st.error(f"Database error: {err}")
      return None

def send_notification(recipient_role, message, username):  
    db = create_connection()  
    if db is None:  
        return False  

    try:  
        notification_doc = {
            "recipient_role": recipient_role,
            "message": message,
            "username": username,
            "is_new": 1,
            "created_at": datetime.now()
        }
        db.notifications.insert_one(notification_doc)
        return True  
    except OperationFailure as err:  
        st.error(f"Database error: {err}")  
        return False  
    
def display_notifications(notifications, username):
    if notifications:
        st.sidebar.write("Notifications:")
        for notification in notifications:
            if notification.get('is_new', 0) == 1:  # Check if the notification is new (MongoDB stores 1/0)
                # Display with a green dot
                st.sidebar.markdown(f"<span style='color: green;'>•</span> {notification.get('message', '')}", unsafe_allow_html=True)
            else:
                st.sidebar.write(notification.get('message', ''))
            
            # Mark as read when clicked (This interaction might need a separate button/logic in Streamlit)
            # For now, this part is commented out as it requires a specific Streamlit callback.
            # if st.sidebar.button(f"Mark as Read - {notification.get('id', '')}"):
            #     mark_notification_as_read(notification.get('_id')) # Use _id for MongoDB
            #     st.experimental_rerun()  # Refresh the app to show updated notifications

    else:
        st.sidebar.write("No notifications available.")

def mark_notification_as_read(notification_id):
    db = create_connection()
    if db is None:
        return False

    try:
        # Convert notification_id to ObjectId if it's a string
        if isinstance(notification_id, str):
            notification_id = ObjectId(notification_id)

        result = db.notifications.update_one({"_id": notification_id}, {"$set": {"is_new": 0}})
        return result.modified_count > 0
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False

def clear_all_notifications(username):
    db = create_connection()
    if db is None:
        return False

    try:
        result = db.notifications.update_many({"username": username}, {"$set": {"is_new": 0}})
        return result.modified_count > 0
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False
  
def get_notifications(recipient_role, username=None):
    db = create_connection()
    if db is None:
        return []

    try:
        query_filter = {"recipient_role": recipient_role}
        if username:
            query_filter["username"] = username
        
        notifications_cursor = db.notifications.find(query_filter).sort("created_at", -1)
        notifications = list(notifications_cursor)
        # Convert ObjectId to string for 'id' if needed for display
        for notification in notifications:
            notification['id'] = str(notification['_id'])
        return notifications  
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return []

def clear_notifications(recipient_role, username):
    db = create_connection()
    if db is None:
        return False

    try:
        result = db.notifications.delete_many({"recipient_role": recipient_role, "username": username})
        print(f"Deleted {result.deleted_count} notifications for {recipient_role} and user {username}.")  # Debugging line
        return result.deleted_count > 0
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False

# New function to update question bank
def update_question_bank(qb_id, new_questions):
    db = create_connection()
    if db is None:
        return False

    try:
        # Convert qb_id to ObjectId if it's a string
        if isinstance(qb_id, str):
            qb_id = ObjectId(qb_id)

        result = db.question_banks.update_one({"_id": qb_id}, {"$set": {"questions": new_questions}})
        return result.modified_count > 0
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False

def update_learning_plan_status(qb_id, username, new_status):   
  db = create_connection()   
  if db is None:   
    return False   
   
  try:
    # Convert qb_id to ObjectId if it's a string
    if isinstance(qb_id, str):
        qb_id_obj = ObjectId(qb_id)
    else:
        qb_id_obj = qb_id

    qb_info = db.question_banks.find_one({"_id": qb_id_obj}, {"technology": 1})
   
    if qb_info and 'technology' in qb_info:   
        technology = qb_info['technology']   
        
        # In MongoDB, you don't alter tables. You just update the document.
        # If 'status' field doesn't exist, it will be added.
        result = db.learning_plans.update_one(
            {"username": username, "question_bank_id": qb_id_obj}, # Use qb_id for unique identification
            {"$set": {"status": new_status, "updated_at": datetime.now()}}, # Add updated_at
            upsert=True # Create if not exists
        )
        return result.modified_count > 0 or result.upserted_id is not None
    else:   
        return False
  except OperationFailure as err:
    st.error(f"Database error: {err}")
    return False

  
def employee_dashboard(username):    
    st.header(f"Welcome, {username}! 🚀") 
    notifications = get_notifications("employee", username)  # Pass the username here  # Pass the username here
    
    if 'chat_history' not in st.session_state:
        st.session_state.chat_history = []  # Initialize chat history 
    user_id = st.session_state.get('user', {}).get('_id', "unique_user_id")  # Use MongoDB _id as user_id

    # Initialize unique chat history for each user
    if user_id and f"chat_history_{user_id}" not in st.session_state:
        st.session_state[f"chat_history_{user_id}"] = []
    with st.sidebar:
        # Create a sidebar for navigation using option_menu
        selected_tab = option_menu(
            menu_title="Employee Dashboard",  # required
            options=["Your Learning Plan", 
                     "Prepare from Generated Questions", 
                     "Take Assessment", 
                     "Your Progress", 
                     "Prepare from Material", 
                     "Discussion Forum", 
                     "Resources",
                     "Chatbot",
                     "Feedback"],  # required
            icons=["book", "pencil", "check-circle", "bar-chart", "file-earmark-text", "chat", "folder", "file-earmark-text"],  # optional
            menu_icon="cast",  # optional
            default_index=0,  # optional
            orientation="vertical",
        )

    # Display the selected tab content
    if selected_tab == "Your Learning Plan":
        st.subheader("Your Learning Plan 📋")
        question_banks = get_all_question_banks()    
        if question_banks:    
            # Create a container for search and dropdown
            search_col1, search_col2 = st.columns([1, 2])
            
            with search_col1:
                # Search input for filtering question banks by ID
                search_id = st.text_input(
                    "Search by ID",
                    key="learning_plan_search_id",
                    placeholder="Enter ID...",
                    help="Enter a question bank ID to quickly find it"
                )
            
            # Format options for the dropdown
            dropdown_options = [(str(qb['_id']), f"{qb['technology']} - {qb['difficulty']}") for qb in question_banks]
            
            # Filter dropdown options if ID is entered
            if search_id:
                try:
                    search_object_id = ObjectId(search_id)
                    dropdown_options = [(str(qb['_id']), f"{qb['technology']} - {qb['difficulty']}")
                                     for qb in question_banks if qb['_id'] == search_object_id]
                    if not dropdown_options:
                        st.warning(f"No question bank found with ID: {search_id}")
                except Exception:
                    st.error("Please enter a valid ID")
            
            with search_col2:
                selected_qb = st.selectbox(    
                    "Select Question Bank for Learning Plan",    
                    options=dropdown_options,    
                    format_func=lambda x: f"ID: {x[0]} - {x[1]}",    
                    key="learning_plan_qb_select",
                    help="Select a question bank from the dropdown or use the ID search to filter"
                )    

            if selected_qb:    
                qb_id_str, _ = selected_qb    
                qb_id = ObjectId(qb_id_str) # Convert to ObjectId for functions that expect it
                learning_plan = prepare_learning_plan(qb_id, username)    
                if learning_plan:    
                    st.write(f"Technology: {', '.join(learning_plan['technology']) if isinstance(learning_plan['technology'], list) else learning_plan['technology']}")    
                    st.write(f"Start Date: {learning_plan['start_date']}")    
                    st.write(f"End Date: {learning_plan['end_date']}")      
                    st.write(f"Estimated Time to Complete: {learning_plan['estimated_time']} hours")    
                    
                    # Update status feature    
                    new_status = st.selectbox("Update Status", ["In Progress", "Completed"], index=0 if learning_plan.get('status') == "In Progress" else 1) # Set default based on current status
                    if st.button("Update Status"):    
                        if update_learning_plan_status(qb_id, username, new_status):    
                            st.success("Status updated successfully!")    
                            # Refresh learning plan details after update
                            st.experimental_rerun()
                        else:    
                            st.error("Failed to update status")    
                    
                    # Generate next learning plan feature    
                    if st.button("Generate Next Learning Plan"):    
                        next_qb_id = get_next_question_bank_id(qb_id)    
                        if next_qb_id:    
                            next_learning_plan = prepare_learning_plan(next_qb_id, username)    
                            if next_learning_plan:    
                                st.write(f"Next Technology: {', '.join(next_learning_plan['technology']) if isinstance(next_learning_plan['technology'], list) else next_learning_plan['technology']}")    
                                st.write(f"Next Start Date: {next_learning_plan['start_date']}")    
                                st.write(f"Next End Date: {next_learning_plan['end_date']}")    
                                st.write(f"Next Status: {next_learning_plan['status']}")    
                                st.write(f"Next Estimated Time to Complete: {next_learning_plan['estimated_time']} hours")    
                            else:    
                                st.error("Failed to generate next learning plan")    
                        else:    
                            st.error("No next question bank available")    
                else:    
                    st.warning("No learning plan assigned yet.")    
        else:    
            st.info("No question banks available yet.")   

    elif selected_tab == "Prepare from Generated Questions":
        st.subheader("Prepare from Generated Questions 📖")    
        question_banks = get_all_question_banks()    
        if question_banks:    
            # Create a container for search and dropdown
            search_col1, search_col2 = st.columns([1, 2])
            
            with search_col1:
                # Search input for filtering question banks by ID
                search_id = st.text_input(
                    "Search by ID",
                    key="qb_search_id",
                    placeholder="Enter ID...",
                    help="Enter a question bank ID to quickly find it"
                )
            
            # Format options for the dropdown
            dropdown_options = [(str(qb['_id']), f"{qb['technology']} - {qb['difficulty']}") for qb in question_banks]
            
            # Filter dropdown options if ID is entered
            if search_id:
                try:
                    search_object_id = ObjectId(search_id)
                    dropdown_options = [(str(qb['_id']), f"{qb['technology']} - {qb['difficulty']}")
                                        for qb in question_banks if qb['_id'] == search_object_id]
                    if not dropdown_options:
                        st.warning(f"No question bank found with ID: {search_id}")
                except Exception:
                    st.error("Please enter a valid ID")
            
            with search_col2:
                selected_qb = st.selectbox(
                    "Select Question Bank", 
                    options=dropdown_options,
                    format_func=lambda x: f"ID: {x[0]} - {x[1]}", 
                    key="prepare_qb_select",
                    help="Select a question bank from the dropdown or use the ID search to filter"
                )    
            
            if selected_qb:    
                qb_id_str, _ = selected_qb    
                qb_id = ObjectId(qb_id_str) # Convert to ObjectId
                qb_details = next((qb for qb in question_banks if qb['_id'] == qb_id), None)    
                if qb_details:    
                    # Display metadata in columns
                    col1, col2, col3 = st.columns(3)  # Removed col4 for updated timestamp as it's not directly displayed
                    with col1:
                        st.info(f"**ID:** {qb_details['_id']}")
                    with col2:
                        st.info(f"**Technology:** {qb_details['technology']}")
                    with col3:
                        st.info(f"**Difficulty:** {qb_details['difficulty']}")
                    

                    # Retrieve and display questions
                    questions = qb_details.get('questions', '').split('\n')    
                    if questions:    
                        st.write("Questions:")
                        for i, question in enumerate(questions, 1):
                            st.write(f"Q{i}: {question}")
                    else:
                        st.warning("No questions available in this question bank.")
                    
                    # Language selection for translation
                    st.subheader("Select Language for Translation")
                    languages = ["English", "Hindi", "Tamil", "Telugu", "Bengali", "Kannada", "Malayalam", 
                                "Spanish", "French", "German", "Chinese", "Japanese", "Korean"]
                    selected_language = st.selectbox("Choose a language", languages)

                    if st.button("Translate Questions"):
                        if questions:
                            translated_questions = []
                            for question in questions:
                                try:
                                    # Use the globally imported GoogleTranslator
                                    translator_instance = GoogleTranslator(source='auto', target=selected_language.lower())
                                    
                                    # Handle language code mapping for certain languages
                                    language_code = selected_language.lower()
                                    # Map language names to their codes when needed
                                    language_mapping = {
                                        "chinese": "zh",
                                        "korean": "ko",
                                        "japanese": "ja"
                                    }
                                    if language_code in language_mapping:
                                        language_code = language_mapping[language_code]
                                        
                                    # Translate using deep-translator
                                    translated = GoogleTranslator(source='auto', target=language_code).translate(question)
                                    translated_questions.append(translated)
                                except Exception as e:
                                    st.error(f"Translation error: {e}")
                                    translated_questions.append(f"Error translating: {question}")
                            
                            st.write("Translated Questions:")
                            for i, question in enumerate(translated_questions, 1):
                                st.write(f"{i}. {question}")
                        else:
                            st.error("No questions available to translate.")



    elif selected_tab == "Take Assessment":
        st.subheader("Take Assessment ✍️")    
        question_banks = get_all_question_banks()    
        if question_banks:    
            # Create a container for search and dropdown
            search_col1, search_col2 = st.columns([1, 2])
            
            with search_col1:
                # Search input for filtering question banks by ID
                search_id = st.text_input(
                    "Search by ID",
                    key="assessment_search_id",
                    placeholder="Enter ID...",
                    help="Enter a question bank ID to quickly find it"
                )
            
            # Format options for the dropdown
            dropdown_options = [(str(qb['_id']), f"{qb['technology']} - {qb['difficulty']}") for qb in question_banks]
            
            # Filter dropdown options if ID is entered
            if search_id:
                try:
                    search_object_id = ObjectId(search_id)
                    dropdown_options = [(str(qb['_id']), f"{qb['technology']} - {qb['difficulty']}")
                                     for qb in question_banks if qb['_id'] == search_object_id]
                    if not dropdown_options:
                        st.warning(f"No question bank found with ID: {search_id}")
                except Exception:
                    st.error("Please enter a valid numeric ID")
            
            with search_col2:
                selected_qb = st.selectbox(    
                    "Select Question Bank",    
                    options=dropdown_options,    
                    format_func=lambda x: f"ID: {x[0]} - {x[1]}",    
                    key="take_assessment_qb_select",
                    help="Select a question bank from the dropdown or use the ID search to filter"
                )    

            if selected_qb:    
                qb_id_str, _ = selected_qb    
                qb_id = ObjectId(qb_id_str) # Convert to ObjectId
                qb_details = next((qb for qb in question_banks if qb['_id'] == qb_id), None)    
                if qb_details:    
                    # Display metadata in columns for better organization
                    col1, col2 = st.columns(2)
                    with col1:
                        st.info(f"**Technology:** {qb_details['technology']}")    
                    with col2:
                        st.info(f"**Difficulty:** {qb_details['difficulty']}")    

                    questions = qb_details.get('questions', '').split('\n')    
                    correct_answers = get_correct_answers(qb_id)    
                    user_answers = []    
                    
                    question_type = qb_details.get('question_type')
                    options_from_qb = qb_details.get('options', '').split('|||') # Split the options string

                    # Assessment questions and answers
                    for i, question in enumerate(questions):
                        if not question.strip(): # Skip empty questions
                            continue

                        st.write(f"**Q{i+1}:** {question.strip()}")
                        
                        if question_type == "multiple-choice":
                            if i < len(options_from_qb):
                                current_options = options_from_qb[i].split('###')
                                if current_options:
                                    # Use a unique key for each radio button group
                                    answer = st.radio("Select an option", current_options, key=f"question_{i}")
                                    user_answers.append(answer)
                                else:
                                    user_answers.append("")
                            else:
                                user_answers.append("")
                        elif question_type == "fill-in-the-blank":
                            answer = st.text_input("Enter your answer", key=f"question_{i}")
                            user_answers.append(answer)
                        elif question_type == "subjective":
                            answer = st.text_area("Enter your answer", key=f"question_{i}")
                            user_answers.append(answer)
                        else:
                            user_answers.append("") # Default if type is unknown or missing

                    if st.button("Submit"):    
                        score = 0    
                        total_questions_answered = 0
                        for i, user_answer in enumerate(user_answers):
                            if i < len(correct_answers) and user_answer.strip().lower() == correct_answers[i].strip().lower():    
                                st.success(f"Q{i+1}: Correct!")    
                                score += 1    
                            elif i < len(correct_answers):
                                st.error(f"Q{i+1}: Incorrect. Correct answer: {correct_answers[i].strip()}")    
                            total_questions_answered += 1
                        
                        # Display final score in a more prominent way
                        st.markdown(f"### Final Score: {score}/{total_questions_answered}")    
                        save_assessment_result(st.session_state.user['username'], qb_id, score)    
                else:    
                    st.error("Failed to retrieve question bank details")    
        else:    
            st.info("No question banks available yet.")

    elif selected_tab == "Your Progress":
        st.subheader("Your Progress 📊")    
        completed_assessments = get_completed_assessments(username)    
        
        if completed_assessments:
            # Create DataFrame for better display
            df = pd.DataFrame(completed_assessments)
            
            # Format completion date
            df['completed_at'] = pd.to_datetime(df['completed_at']).dt.strftime('%Y-%m-%d %H:%M')
            
            # Add filter options
            col1, col2, col3 = st.columns(3)
            with col1:
                tech_filter = st.multiselect(
                    "Filter by Technology",
                    options=sorted(df['technology'].unique()),
                    default=[],
                    help="Select one or more technologies to filter"
                )
            
            with col2:
                difficulty_filter = st.multiselect(
                    "Filter by Difficulty",
                    options=sorted(df['difficulty'].unique()),
                    default=[],
                    help="Select one or more difficulty levels to filter"
                )
                
            with col3:
                question_type_filter = st.multiselect(
                    "Filter by Question Type",
                    options=sorted(df['question_type'].unique()),
                    default=[],
                    help="Select one or more question types to filter"
                )
            
            # Apply filters
            filtered_df = df.copy()
            if tech_filter:
                filtered_df = filtered_df[filtered_df['technology'].isin(tech_filter)]
            if difficulty_filter:
                filtered_df = filtered_df[filtered_df['difficulty'].isin(difficulty_filter)]
            if question_type_filter:
                filtered_df = filtered_df[filtered_df['question_type'].isin(question_type_filter)]
            
            # Display statistics
            st.subheader("Summary Statistics")
            stats_col1, stats_col2, stats_col3, stats_col4 = st.columns(4)
            
            with stats_col1:
                st.metric(
                    "Total Assessments",
                    len(filtered_df),
                    help="Total number of completed assessments"
                )
            
            with stats_col2:
                avg_score = filtered_df['percentage'].mean()
                st.metric(
                    "Average Score",
                    f"{avg_score:.1f}%",
                    help="Average score across all completed assessments"
                )
            
            with stats_col3:
                best_score = filtered_df['percentage'].max()
                st.metric(
                    "Best Score",
                    f"{best_score:.1f}%",
                    help="Your highest score across all assessments"
                )
                
            with stats_col4:
                recent_trend = filtered_df.head(3)['percentage'].mean()
                st.metric(
                    "Recent Performance",
                    f"{recent_trend:.1f}%",
                    help="Average score of your last 3 assessments"
                )
            
            # Display the table with proper formatting
            st.subheader("Detailed Progress")
            
            # Create a styled dataframe
            display_df = filtered_df[[
                'id', 'technology', 'difficulty', 'question_type', 'score', 
                'total_questions', 'percentage', 'completed_at'
            ]].rename(columns={
                'id': 'Assessment ID',
                'technology': 'Technology',
                'difficulty': 'Difficulty',
                'question_type': 'Question Type',
                'score': 'Score',
                'total_questions': 'Total Questions',
                'percentage': 'Percentage (%)',
                'completed_at': 'Completed At'
            })
            
            # Apply styling based on percentage
            def color_percentage(val):
                try:
                    val_num = float(val)
                    if val_num >= 80:
                        return 'background-color: #90EE90'  # Light green
                    elif val_num >= 60:
                        return 'background-color: #FFD700'  # Gold
                    return 'background-color: #FFB6C1'  # Light red
                except:
                    return ''
            
            styled_df = display_df.style.applymap(
                color_percentage, 
                subset=['Percentage (%)']
            )
            
            st.dataframe(
                styled_df,
                use_container_width=True,
                hide_index=True
            )
            
            # Add chart to visualize progress over time
            st.subheader("Progress Over Time")
            fig = px.line(
                display_df.sort_values('Completed At'), 
                x='Completed At', 
                y='Percentage (%)',
                title='Score Trend',
                markers=True
            )
            st.plotly_chart(fig, use_container_width=True)
            
            # Add download button for the progress report
            csv = display_df.to_csv(index=False)
            st.download_button(
                label="Download Progress Report",
                data=csv,
                file_name="assessment_progress.csv",
                mime="text/csv",
                help="Download your progress report as a CSV file"
            )
            
        else:    
            st.info("You haven't completed any assessments yet.")

    elif selected_tab == "Prepare from Material":
            st.subheader("Prepare from Material 📚")    
            curricula = get_all_curricula()    
            if curricula:
                # Add a "None" placeholder as the first option
                curriculum_options = ["None"] + [c['technology'] for c in curricula]
                
                selected_curriculum = st.selectbox(    
                    "Select Curriculum to Prepare from",    
                    options=curriculum_options,    
                    key="prepare_curriculum_select"    
                )    

                # Only display content if a valid curriculum is selected
                if selected_curriculum and selected_curriculum != "None":    
                    curriculum_content = get_curriculum_text(selected_curriculum)    
                    if curriculum_content:
                        # curriculum_content is already a string from MongoDB
                        decoded_content = curriculum_content
                        
                        # Clean up special characters
                        decoded_content = decoded_content.replace('\xef\x82\xb7', '•').replace('\xef\x80\xa0', '')  
                        
                        st.write("Curriculum Content:")  
                        st.write(decoded_content)  

                        st.subheader("Select Language for Translation")
                        languages = ["English", "Spanish", "French", "German", "Chinese", "Japanese", "Korean","Kannada","Malayalam", "Hindi", "Tamil", "Telugu", "Bengali"]
                        selected_language = st.selectbox("Choose a language", languages)

                        if st.button("Translate Curriculum"):
                            if decoded_content:
                                try:
                                    # Using deep-translator for more reliable translation
                                    # Removed the redundant import here
                                    translator_instance = GoogleTranslator(source='auto', target=selected_language.lower())
                                    
                                    # For long texts, we might need to split and translate in chunks
                                    max_chunk = 5000  # Google Translate has a limit
                                    
                                    if len(decoded_content) > max_chunk:
                                        # Split into chunks and translate separately
                                        chunks = [decoded_content[i:i+max_chunk] for i in range(0, len(decoded_content), max_chunk)]
                                        translated_chunks = [translator_instance.translate(chunk) for chunk in chunks]
                                        translated_content = ''.join(translated_chunks)
                                    else:
                                        translated_content = translator_instance.translate(decoded_content)
                                        
                                    st.write("Translated Curriculum Content:")
                                    st.write(translated_content)
                                except Exception as e:
                                    st.error(f"Translation failed: {e}")
                            else:
                                st.error("No content available to translate.")

                        # Rest of your document download code...
                        if "curriculum_downloaded" not in st.session_state:
                            st.session_state.curriculum_downloaded = False

                        if not st.session_state.curriculum_downloaded:
                            if st.button("Download Curriculum"):  
                                import docx
                                import io
                                doc = docx.Document()  
                                doc.add_paragraph(decoded_content)  
                                buffer = io.BytesIO()  
                                doc.save(buffer)  
                                buffer.seek(0)  
                                st.download_button(  
                                    label="Download Curriculum",  
                                    data=buffer.getvalue(),  
                                    file_name=f"curriculum_{selected_curriculum}.docx",  
                                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"  
                                )  
                                st.session_state.curriculum_downloaded = True
                        else:
                            st.info("Curriculum has already been downloaded.")
                    else:    
                        st.error("Failed to retrieve curriculum content")    
            else:    
                st.info("No curricula available yet.")


    elif selected_tab == "Discussion Forum":
        discussion_forum()

    elif selected_tab == "Resources":
        st.subheader("Resources 🌐")
        st.write("Here you can find various resources to help you in your learning journey.")

        # Search bar for technology
        search_skill = st.text_input("Search for a technology or skill:")

        if search_skill:
            resources_html = generate_resources(search_skill)
            st.markdown(resources_html, unsafe_allow_html=True)
        else:
            print("Please enter Technology")
    elif selected_tab == "Chatbot":
        # Display chatbot interface at the top
        st.subheader("Chat with the AI Trainer 🤖")

        # Initialize message state if it doesn't exist
        if "msg" not in st.session_state:
            st.session_state.msg = ""
            
        # Create a container for the chat interface
        chat_container = st.container()

        # Define avatars
        user_avatar = "https://static.vecteezy.com/system/resources/previews/009/664/418/non_2x/people-user-team-transparent-free-png.png"
        ai_avatar = "https://thumbs.dreamstime.com/b/chatbot-logo-messenger-ai-robot-icon-vector-illustration-277900892.jpg"

        def clear_text():
            st.session_state.msg = st.session_state.user_input
            st.session_state.user_input = ""

        with chat_container:
            # Display chat messages specific to this user
            for chat in st.session_state[f"chat_history_{user_id}"]:
                if chat['role'] == 'assistant':
                    # Chatbot message with avatar
                    st.markdown(
                        f"<div style='display: flex; align-items: center; margin: 5px 0;'>"
                        f"<img src='{ai_avatar}' style='width: 40px; height: 40px; border-radius: 50%; margin-right: 10px;'>"
                        f"<div style='background-color: #e1ffc7; padding: 10px; border-radius: 10px; max-width: 80%;'>"
                        f"<strong>AI:</strong> {chat['content']}</div></div>",
                        unsafe_allow_html=True
                    )
                else:
                    # User message with avatar
                    st.markdown(
                        f"<div style='display: flex; align-items: center; margin: 5px 0; justify-content: flex-end;'>"
                        f"<div style='background-color: #dcf8c6; padding: 10px; border-radius: 10px; max-width: 80%; margin-left: auto;'>"
                        f"<strong>You:</strong> {chat['content']}</div>"
                        f"<img src='{user_avatar}' style='width: 40px; height: 40px; border-radius: 50%; margin-left: 10px;'>"
                        f"</div>", 
                        unsafe_allow_html=True
                    )

            # Input field for user to enter a prompt
            st.text_input("Type your message here...", key="user_input", placeholder="Type a message...", on_change=clear_text)

            if st.session_state.msg:  # Only process if there's a message
                # Append user input to this user's chat history
                st.session_state[f"chat_history_{user_id}"].append({"role": "user", "content": st.session_state.msg})

                try:
                    # Generate AI response
                    prompt = f"You are an AI assistant for trainers. Respond to the following message: {st.session_state.msg}"
                    response = model.generate_content(prompt)

                    # Handle the response based on the model's response format
                    if hasattr(response, 'parts'):
                        ai_response = ''.join(part.text for part in response.parts)
                    else:
                        ai_response = response.candidates[0].content.parts[0].text

                    # Append AI response to this user's chat history
                    st.session_state[f"chat_history_{user_id}"].append({"role": "assistant", "content": ai_response})
                except Exception as e:
                    st.error(f"Error generating response: {str(e)}")
                    ai_response = "I apologize, but I encountered an error. Please try again."
                    st.session_state[f"chat_history_{user_id}"].append({"role": "assistant", "content": ai_response})

                # Clear the message state
                st.session_state.msg = ""
                # Rerun the app to display the new messages
                #st.experimental_rerun()
                st.rerun()

    
    elif selected_tab == "Feedback":
        st.subheader("Submit Feedback ✍🏻")
        
        # Select the type of feedback
        feedback_type = st.radio("Select Feedback Type:", 
                                options=["User  Experience", "Question Bank Feedback", "Assessment Feedback"])
        
        # Select a question bank for feedback if applicable
        question_banks = get_all_question_banks()
        if feedback_type in ["Question Bank Feedback", "Assessment Feedback"] and question_banks:
            selected_qb = st.selectbox(
                "Select Question Bank for Feedback",
                options=[(str(qb['_id']), f"{qb['technology']} - {qb['difficulty']}") for qb in question_banks],
                format_func=lambda x: f"ID: {x[0]} - {x[1]}",
                key="feedback_qb_select"
            )
        else:
            selected_qb = None  # No question bank needed for User Experience feedback

        feedback_text = st.text_area("Submit your feedback:", height=100)
        rating = st.slider("Rate your experience (1-5):", 1, 5, 3)

        if st.button("Submit Feedback"):
            if feedback_type == "User  Experience":
                question_bank_id = None  # No question bank for user experience feedback
            else:
                question_bank_id = selected_qb[0] if selected_qb else None

            if submit_feedback(username, question_bank_id, feedback_text, rating, feedback_type):
                st.success("Feedback submitted successfully!")
            else:
                st.error("Failed to submit feedback.")

    notifications = get_notifications("employee", username)  # Pass the username here
    display_notifications(notifications, username)  # Display notifications in the sidebar
    if notifications:  
        st.sidebar.write("Notifications:")  
        for notification in notifications:  
            st.sidebar.write(notification['message'])  
    else:  
        st.sidebar.write("No notifications available.")

# Main function for the discussion forum  
def discussion_forum():  
    st.title("Discussion Forum 💬")  

    # Add a refresh icon beside the title  
    refresh_icon = '🔄'  
    refresh_button = st.button(refresh_icon, key="refresh_forum_button")  

    # Section to post a new message  
    new_message = st.text_area("Post a new message", height=100)  
    if st.button("Post Message", key="post_message_button"):  
        if new_message and st.session_state.user['username']:  
            if save_message(new_message, st.session_state.user['username']):  
                st.success("Message posted successfully!")  
            else:  
                st.error("Failed to post message.")  
        else:  
            st.error("Message and username cannot be empty.")  

    # Display all messages in reverse order to show the newest first  
    messages = get_messages()  
    if messages:  
        for message in reversed(messages):  # Reverse the order of messages
            # Create a container for each message
            message_container = st.container()
            with message_container:
                # Display the message with the username
                if message.get('username') == st.session_state.user['username']:  # Check if the message is from the user
                    # User message on the right
                    st.markdown(
                        f"<div style='display: flex; align-items: center; justify-content: flex-end; margin: 5px 0;'>"
                        f"<div style='background-color: #dcf8c6; padding: 10px; border-radius: 10px; max-width: 80%; margin-left: auto;'>"
                        f"<strong>You:</strong> {message.get('message', '')}</div>"
                        f"</div>", 
                        unsafe_allow_html=True
                    )
                else:
                    # Other user's message on the left
                    st.markdown(
                        f"<div style='display: flex; align-items: center; margin: 5px 0;'>"
                        f"<div style='background-color: #e1ffc7; padding: 10px; border-radius: 10px; max-width: 80%;'>"
                        f"<strong>{message.get('username', 'Unknown')}:</strong> {message.get('message', '')}</div></div>",
                        unsafe_allow_html=True
                    )

                # Fetch and display replies  
                replies = get_replies(message['_id'])  # Use MongoDB _id
                if replies:  
                    for reply in replies:  
                        # Display each reply
                        if reply.get('username') == st.session_state.user['username']:  # Check if the reply is from the user
                            # User reply on the right
                            st.markdown(
                                f"<div style='display: flex; align-items: center; justify-content: flex-end; margin: 5px 0;'>"
                                f"<div style='background-color: #dcf8c6; padding: 10px; border-radius: 10px; max-width: 80%; margin-left: auto;'>"
                                f"<strong>You:</strong> {reply.get('reply_message', '')}</div>"
                                f"</div>", 
                                unsafe_allow_html=True
                            )
                        else:
                            # Other user's reply on the left
                            st.markdown(
                                f"<div style='display: flex; align-items: center; margin: 5px 0;'>"
                                f"<div style='background-color: #e1ffc7; padding: 10px; border-radius: 10px; max-width: 80%;'>"
                                f"<strong>{reply.get('username', 'Unknown')}:</strong> {reply.get('reply_message', '')}</div></div>",
                                unsafe_allow_html=True
                            )  

                else:  
                    st.write("  No replies yet.")  

            # Right column for reply option  
            reply_message = st.text_area("Reply", key=f"reply_area_{message['_id']}")  
            if st.button("Post Reply", key=f"reply_button_{message['_id']}"):  
                if reply_message:  
                    if save_reply(message['_id'], reply_message, st.session_state.user['username']):  
                        st.success("Reply posted successfully!")  
                        st.experimental_rerun() # Refresh to show new reply
                    else:  
                        st.error("Failed to post reply.")  
                else:  
                    st.error("Reply cannot be empty.")  

    else:  
        st.info("No messages available.")  

    # Refresh the forum if the refresh button is clicked  
    if refresh_button:  
        st.session_state.needs_refresh = True  

    if 'needs_refresh' in st.session_state and st.session_state.needs_refresh:  
        st.session_state.needs_refresh = False  
        # Perform any necessary actions here, such as updating the messages or replies  
        st.write("Forum refreshed!")  

def generate_resources(skill):  
   genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))  # Use os.getenv for API key
   model = genai.GenerativeModel("gemini-1.5-flash")  
   response = model.generate_content(f"Provide a list of resources (articles, videos, etc.) related to the skill '{skill}'. Format the response in an HTML table with the following CSS: table {{ width: 100%; border-collapse: collapse; font-family: sans-serif; }} th, td {{ padding: 10px; text-align: left; border-bottom: 1px solid #ddd; }} th {{ background-color: #3498db; color: white; }}")  
   return response.text.strip()

def get_replies(message_id):
    db = create_connection()
    if db is None:
        return []

    try:
        replies_cursor = db.replies.find({"message_id": message_id}, {"reply_message": 1, "username": 1, "_id": 0})
        replies = list(replies_cursor)
        return replies
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return []

def get_resources():  
    db = create_connection()  
    if db is None:
        return []
    try:
        resources_cursor = db.resources.find({})  
        resources = list(resources_cursor)  
        return resources  
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return []

def filter_resources(category):  
    db = create_connection()  
    if db is None:
        return []
    try:
        filtered_resources_cursor = db.resources.find({"category": category})  
        filtered_resources = list(filtered_resources_cursor)  
        return filtered_resources  
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return []

def add_widget(widget):  
    db = create_connection()  
    if db is None:
        return False
    try:
        db.widgets.insert_one({"widget": widget})  
        return True
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False

def get_widgets():  
    db = create_connection()  
    if db is None:
        return []
    try:
        widgets_cursor = db.widgets.find({})  
        widgets = list(widgets_cursor)  
        return widgets  
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return []

def widget_settings(widget_id):  
    db = create_connection()  
    if db is None:
        return None
    try:
        # Assuming widget_id is the MongoDB _id
        widget = db.widgets.find_one({"_id": ObjectId(widget_id)})  
        return widget  
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None

def change_layout(layout):  
    db = create_connection()  
    if db is None:
        return False
    try:
        result = db.users.update_one({"username": st.session_state.user['username']}, {"$set": {"layout": layout}})  
        return result.modified_count > 0
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False

# Database functions  
def save_message(message, username):
    db = create_connection()
    if db is None:
        st.error("Failed to connect to the database.")
        return False

    try:
        message_doc = {
            "username": username,
            "message": message,
            "created_at": datetime.now()
        }
        db.messages.insert_one(message_doc)
        return True
    except OperationFailure as err:
        st.error(f"Database error: {err}")  # Show the specific database error
        return False
  
def get_messages():  
    db = create_connection()  
    if db is None:
        return []
    try:
        messages_cursor = db.messages.find({})  
        messages = list(messages_cursor)  
        return messages
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return []
  
def save_reply(message_id, reply_message, username):
    db = create_connection()
    if db is None:
        return False

    try:
        reply_doc = {
            "message_id": message_id, # This should be the ObjectId of the parent message
            "reply_message": reply_message,
            "username": username,
            "created_at": datetime.now()
        }
        db.replies.insert_one(reply_doc)
        return True
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False 
  
def search_messages(search_query):  
    db = create_connection()  
    if db is None:
        return []
    try:
        # Use regex for case-insensitive search
        search_results_cursor = db.messages.find({"message": {"$regex": search_query, "$options": "i"}})  
        search_results = list(search_results_cursor)  
        return search_results
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return []

def get_generated_questions():
    db = create_connection()
    if db is None:
        return None

    try:
        result = db.generated_question_files.find_one({}, {"questions": 1, "options": 1, "correct_answers": 1, "_id": 0})
        if result:
            return {
                'questions': result.get('questions', ''),
                'options': result.get('options', ''),
                'correct_answers': result.get('correct_answers', '')
            }
        else:
            return None
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return None

def get_assessment_results(username):  
   db = create_connection()  
   if db is None:  
      return None  
  
   try:  
      results_cursor = db.assessments.find({"username": username})  
      results = list(results_cursor)  
      # Convert ObjectId to string for question_bank_id if it exists
      for result in results:
          if 'question_bank_id' in result and isinstance(result['question_bank_id'], ObjectId):
              result['question_bank_id'] = str(result['question_bank_id'])
      return results  
   except OperationFailure as err:  
      st.error(f"Database error: {err}")  
      return None

def get_next_question_bank_id(qb_id):   
  db = create_connection()   
  if db is None:   
    return None   
   
  try:
    # Convert qb_id to ObjectId if it's a string
    if isinstance(qb_id, str):
        qb_id_obj = ObjectId(qb_id)
    else:
        qb_id_obj = qb_id

    # Find the next question bank by _id (which is usually ordered chronologically)
    result = db.question_banks.find({"_id": {"$gt": qb_id_obj}}).sort("_id", 1).limit(1)
    result_doc = next(result, None)
    
    if result_doc:   
      return str(result_doc['_id']) # Return as string
    else:   
      return None
  except OperationFailure as err:
    st.error(f"Database error: {err}")
    return None

def prepare_learning_plan(qb_id, username):   
  db = create_connection()   
  if db is None:   
    return None   
   
  try:
    # Convert qb_id to ObjectId if it's a string
    if isinstance(qb_id, str):
        qb_id_obj = ObjectId(qb_id)
    else:
        qb_id_obj = qb_id

    qb_details = db.question_banks.find_one({"_id": qb_id_obj})
   
    if qb_details:   
        questions = qb_details.get('questions', '').split('\n')   
        difficulty = qb_details.get('difficulty')   
        num_questions = len(questions)   
        estimated_time = calculate_estimated_time(num_questions, difficulty)   
        topics = []   
        # Extract topics from questions (assuming format "Topic: Question")
        for question in questions:   
            if ':' in question:
                topic = question.split(':', 1)[0].strip()   
                if topic not in topics:   
                    topics.append(topic)   
        
        # Get the date when the employee updated the status as 'Completed'   
        completed_date = get_completed_date(username)   
        
        # Set the start date to the day after the completed date   
        start_date = (completed_date + timedelta(days=1)).strftime('%Y-%m-%d')   
        
        # Calculate the estimated end date based on the topics length and other analysis   
        estimated_end_date = calculate_estimated_end_date(topics, estimated_time, start_date)   
        
        learning_plan_data = {   
            'technology': topics,   
            'start_date': start_date,   
            'end_date': estimated_end_date,   
            'status': 'In Progress',   
            'estimated_time': estimated_time,
            'username': username, # Ensure username is part of the learning plan
            'question_bank_id': qb_id_obj # Store ObjectId
        }   
        
        # Upsert the learning plan: update if exists, insert if not
        db.learning_plans.update_one(
            {"username": username, "question_bank_id": qb_id_obj},
            {"$set": learning_plan_data},
            upsert=True
        )

        return learning_plan_data   
    else:   
        return None
  except OperationFailure as err:
    st.error(f"Database error: {err}")
    return None

def get_correct_answers(qb_id):  
   db = create_connection()  
   if db is None:  
      return None  
  
   try:  
      # Convert qb_id to ObjectId if it's a string
      if isinstance(qb_id, str):
          qb_id_obj = ObjectId(qb_id)
      else:
          qb_id_obj = qb_id

      result = db.question_answers.find_one({"question_bank_id": qb_id_obj}, {"answer_data": 1, "_id": 0})
  
      if result and 'answer_data' in result:  
        return result['answer_data'].split('\n')  
      else:  
        return []  
   except OperationFailure as err:  
      st.error(f"Database error: {err}")  
      return None

def save_assessment_result(username, qb_id, score):
    db = create_connection()
    if db is None:
        return False

    try:
        # Convert qb_id to ObjectId if it's a string
        if isinstance(qb_id, str):
            qb_id_obj = ObjectId(qb_id)
        else:
            qb_id_obj = qb_id

        assessment_doc = {
            "username": username,
            "question_bank_id": qb_id_obj, # Store as ObjectId
            "score": score,
            "completed_at": datetime.now()
        }
        db.assessments.insert_one(assessment_doc)
        
        # Send notification based on score
        if score >= 8:
            feedback_message = f"Great job, {username}! Your score of {score}/10 is excellent!"
        elif score >= 5:
            feedback_message = f"Good effort, {username}! Your score of {score}/10 shows progress."
        else:
            feedback_message = f"Keep practicing, {username}. Your score of {score}/10 indicates more review is needed."
        
        send_notification("employee", feedback_message, username)  # Pass the username here
        
        return True
        
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False

def get_previous_learning_plan_end_date(username):   
  db = create_connection()   
  if db is None:   
    return None   
   
  try:
    result = db.learning_plans.find_one(
        {"username": username, "status": "Completed"},
        {"end_date": 1, "_id": 0},
        sort=[("end_date", -1)] # Sort by end_date descending
    )
   
    if result and 'end_date' in result:   
      return datetime.strptime(result['end_date'], '%Y-%m-%d')   
    else:   
      return datetime.now()
  except OperationFailure as err:
    st.error(f"Database error: {err}")
    return datetime.now() # Return current time on error

def get_completed_date(username):   
  db = create_connection()   
  if db is None:   
    return None   
   
  try:
    # In MongoDB, you don't alter tables to add columns. Just ensure the field is present.
    # The 'updated_at' field will be added when a document is updated/inserted with it.
    
    result = db.learning_plans.find_one(
        {"username": username, "status": "Completed"},
        {"updated_at": 1, "_id": 0},
        sort=[("updated_at", -1)] # Sort by updated_at descending
    )
   
    if result and 'updated_at' in result:   
      return result['updated_at'] # updated_at is stored as datetime object
    else:   
      return datetime.now()
  except OperationFailure as err:
    st.error(f"Database error: {err}")
    return datetime.now() # Return current time on error
   
def calculate_estimated_time(num_questions, difficulty):   
  if difficulty == 'Easy':   
    estimated_time = num_questions * 10  # 10 minutes per question   
  elif difficulty == 'Medium':   
    estimated_time = num_questions * 20  # 20 minutes per question   
  elif difficulty == 'Hard':   
    estimated_time = num_questions * 30  # 30 minutes per question   
  return estimated_time
def calculate_estimated_end_date(topics, estimated_time, start_date):   
  # Calculate the estimated end date based on the topics length and other analysis   
  # For example, assume each topic takes 1 day to complete   
  num_days = len(topics)   
  estimated_end_date = (datetime.strptime(start_date, '%Y-%m-%d') + timedelta(days=num_days)).strftime('%Y-%m-%d')   
  return estimated_end_date

import streamlit as st
import requests
from streamlit_option_menu import option_menu
from streamlit_lottie import st_lottie

# Function to load Lottie animation from a URL
def load_lottie_url(url: str):
    try:
        response = requests.get(url)
        if response.status_code == 200:
            return response.json()
        else:
            st.error(f"Failed to load Lottie animation. Status code: {response.status_code}")
            return None
    except Exception as e:
        st.error(f"Error loading Lottie animation: {e}")
        return None

def check_admin_exists():
    """Check if an administrator account already exists"""
    db = create_connection()
    if db is None:
        return False

    try:
        admin_count = db.users.count_documents({"role": "Administrator"})
        return admin_count > 0
    except OperationFailure as err:
        st.error(f"Database error: {err}")
        return False

def main():
    
    if 'user' not in st.session_state:
        st.session_state.user = None

    if st.session_state.user is None:
        # Show the two-column layout only during login/register
        col1, col2 = st.columns(2)

        # Column 1: Embed Lottie animation
        with col1:
            st.components.v1.html(
                """
                <iframe src="https://lottie.host/embed/1b7b20ac-876d-4a6f-82d5-a1b188f88863/6aZt4s4ExJ.json" 
                        width="100%" height="600" frameborder="0" allowfullscreen></iframe>
                """,
                height=600,
            )

        # Column 2: User authentication
        with col2:
            st.title("Automated Question Builder")
            st.write("")
            st.write("")
            st.write("")
            
            selected = option_menu(
                menu_title=None,
                options=["Login", "Register"],
                icons=["person", "person-plus"],
                menu_icon="cast",
                default_index=0,
                orientation="horizontal",
            )

            if selected == "Login":
                username = st.text_input("Username 👤", key="login_username", 
                                      placeholder="Enter your username", help="Your username")
                password = st.text_input("Password 🔑", type="password", key="login_password", 
                                       placeholder="Enter your password", help="Your password")
                if st.button("Login", key="login_button"):
                    user = login_user(username, password)
                    if user:
                        st.session_state.user = user
                        st.success("Logged in successfully!")
                        st.rerun()
                    else:
                        st.error("Invalid username or password")

            elif selected == "Register":
                new_email = st.text_input("Email ✉️ ", key="register_email", 
                                            placeholder="Enter your email")
                new_username = st.text_input("Username 👤", key="register_username", 
                                            placeholder="Choose a username")
                new_password = st.text_input("Password 🔑", type="password", 
                                            key="register_password", 
                                            placeholder="Choose a password")

                # Check if admin exists before showing admin role option
                admin_exists = check_admin_exists()
                if admin_exists:
                    role_options = ["Trainer", "Employee"]
                    role = st.selectbox("Role 👨🏻‍💼", role_options, key="register_role")
                else:
                    role_options = ["Administrator", "Trainer", "Employee"]
                    role = st.selectbox("Role", role_options, key="register_role")
                    if role == "Administrator":
                        st.warning("You are registering as the system administrator. This role can only be assigned once.")

                if st.button("Register", key="register_button"):
                    # Double check admin existence before registration
                    if role == "Administrator" and check_admin_exists():
                        st.error("An administrator account already exists. Please select a different role.")
                    else:
                        if register_user(new_email, new_username, new_password, role):
                            st.success("Registration successful! Please log in.")
                        else:
                            st.error("Registration failed. Username may already exist.")

    else:
        # Single column layout for logged-in users
        
        
        st.sidebar.write(f"Logged in as: {st.session_state.user['username']}")
        if st.sidebar.button("Logout", key="logout_button"):
            st.session_state.user = None
            st.rerun()

        if st.session_state.user['role'] == 'Administrator':
            admin_dashboard()
        elif st.session_state.user['role'] == 'Trainer':
            trainer_dashboard()
        elif st.session_state.user['role'] == 'Employee':
            employee_dashboard(st.session_state.user['username'])

if __name__ == "__main__":
    main()
